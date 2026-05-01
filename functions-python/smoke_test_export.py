"""ローカル品質検証スクリプト: Excel/PPTX を 3 シナリオで生成し、構造を読み戻して検証。

実行: venv/Scripts/python smoke_test_export.py
出力: out_smoke/scenario_{1,2,3}.xlsx + .pptx
"""

from __future__ import annotations

import io
import os
import sys
from pathlib import Path

# Windows コンソールで em-dash 等の Unicode を出力するため utf-8 強制
if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# このスクリプト自身の親ディレクトリ (functions-python) を sys.path へ
sys.path.insert(0, str(Path(__file__).resolve().parent))

from export_excel.builder import build_excel_workbook
from export_pptx.builder import build_pptx_presentation


OUT_DIR = Path(__file__).parent / "out_smoke"
OUT_DIR.mkdir(exist_ok=True)


# ─── 共通ダミーデータ ─────────────────────────
def _dynamic_sheet(rows: int = 5, with_comp: bool = False) -> dict:
    base = [
        {"date": f"2026-04-{d:02d}", "sessions": 100 + d * 10, "totalUsers": 80 + d * 8,
         "screenPageViews": 200 + d * 15, "engagementRate": 0.7, "conversions": d, "channelName": "Direct"}
        for d in range(1, rows + 1)
    ]
    return {
        "visibleColumns": [
            {"key": "date", "label": "日付", "format": "text"},
            {"key": "sessions", "label": "セッション", "format": "number"},
            {"key": "totalUsers", "label": "ユーザー", "format": "number"},
            {"key": "screenPageViews", "label": "PV", "format": "number"},
            {"key": "engagementRate", "label": "ENG率", "format": "percent"},
            {"key": "conversions", "label": "CV", "format": "number"},
        ],
        "rows": base,
        "compRows": [{**r, "sessions": r["sessions"] - 10} for r in base] if with_comp else None,
    }


def _user_journey(num_paths: int = 6) -> dict:
    nodes = [
        {"id": "src1", "name": "オーガニック検索", "type": "source", "value": 3286, "share": 0.441, "change": 0.05},
        {"id": "src2", "name": "ダイレクト", "type": "source", "value": 1973, "share": 0.265, "change": -0.02},
        {"id": "src3", "name": "オーガニックSNS", "type": "source", "value": 1385, "share": 0.186, "change": 0.18},
        {"id": "lp1", "name": "/", "type": "lp", "value": 1293},
        {"id": "lp2", "name": "/archives/2502/long-path-name-here-very-long", "type": "lp", "value": 915},
        {"id": "lp3", "name": "/works/category/web-production", "type": "lp", "value": 254},
        {"id": "cv1", "name": "資料ダウンロード_Webサイト制作の流れ完了", "type": "cv", "value": 25, "share": 0.0034},
        {"id": "cv2", "name": "新規お問い合わせ送信完了", "type": "cv", "value": 14, "share": 0.0019},
        {"id": "cv3", "name": "見積り依頼完了", "type": "cv", "value": 9, "share": 0.0012},
    ]
    story_top3 = [
        {"rank": 1, "title": "自然検索 → トップ → 資料ダウンロード Webサイト制作の流れ完了",
         "sharePct": 8, "cvRate": 2.6, "sessions": 617, "type": "normal",
         "aiComment": "自然検索経由の「資料ダウンロード Webサイト制作の流れ完了」は CV 率 2.6% と健闘していますが、セッション数 8% の主要流入経路としては、より高い CV 率を狙う余地があります。LP 上のメインビジュアルでダウンロード価値を端的に伝える工夫が必要です。"},
        {"rank": 2, "title": "資料ダウンロード Webサイト制作の流れ", "sharePct": 6, "cvRate": 3.3, "sessions": 482,
         "type": "success", "aiComment": "短文 AI コメント"},
        {"rank": 3, "title": "自然検索 → (不明) → 離脱", "sharePct": 6, "cvRate": 0.0, "sessions": 430, "type": "warning",
         "aiComment": "(不明) ページを経由した離脱パターンが顕著。リファラ取得失敗の可能性があるため、計測タグの再確認が必要です。"},
    ]
    detail_paths = [
        {"rank": i, "source": "自然検索" if i % 2 else "ダイレクト",
         "lp": f"/archives/{i}502/very-long-page-path-{i}",
         "middle": "/works/" if i < 4 else "—",
         "result": "資料ダウンロード_Webサイト制作の流れ完了" if i % 3 == 0 else "離脱",
         "sessions": 700 - i * 70, "cvRate": 2.6 - i * 0.3, "change": 0.0}
        for i in range(1, num_paths + 1)
    ]
    return {"totalSessions": 7457, "nodes": nodes, "links": [], "storyTop3": story_top3, "detailPaths": detail_paths}


def _keywords_funnel(with_clusters: bool = True) -> dict:
    layers = ["branded", "pureIntent", "intent", "latent", "noise"]
    funnel = {layer: {
        "layerKey": layer, "labelJa": layer,
        "count": 100 + i * 50, "clicks": 400 - i * 50, "impressions": 10000 + i * 5000,
        "ctr": 4.0 - i * 0.5, "avgPosition": 7.7 + i * 2.5, "estimatedCV": max(0, 4 - i),
        "topKeywords": [f"きーわーど{layer}{n}" for n in range(1, 4)],
        "change": 0.0,
    } for i, layer in enumerate(layers)}
    clusters = []
    if with_clusters:
        names = ["指名検索系", "Web 制作・開発系", "フォント・ファイル形式系", "ユーザー理解系", "デザイン・UI/UX 系", "その他"]
        clusters = [{"id": f"c{i}", "name": n, "centerKeyword": f"中心キーワード{i}",
                     "keywordCount": 50 - i * 5, "clicks": 300 - i * 40, "impressions": 5000 - i * 500}
                    for i, n in enumerate(names, 1)]
    keywords = [
        {"query": f"検索キーワード {i}", "clicks": 50 - i, "impressions": 1000 - i * 50,
         "ctr": 0.05, "position": 8.0 + i, "topPage": f"/works/page-{i}-very-long-url-here",
         "layer": "intent", "ctrLossFlag": i < 5, "ctrLossDelta": -2.5,
         "potentialClicks": 30 - i * 2, "estimatedCV": 0, "change": 0.0}
        for i in range(1, 11)
    ]
    return {
        "metrics": {"keywordCount": 4470, "totalClicks": 800, "totalImpressions": 58800,
                    "avgCTR": 1.36, "avgPosition": 12.12, "estimatedCV": 5},
        "funnel": funnel, "clusters": clusters, "keywords": keywords,
    }


# 短文/長文/極長の AI 分析
SHORT_AI = "簡潔な AI 分析。"
MID_AI = "当期のデータでは大きな変動が観測されています。\n・CV 率は前月比 +0.3pt\n・主要な改善余地は LP の遷移率"
LONG_AI = "\n".join([
    "当期のデータによると、総セッション数は87,330回、総コンバージョン数は660件でした。月別に見ると、ユーザー数やセッション数は比較的安定していますが、コンバージョン数には月ごとの変動が見られます。特に2025年12月と2026年01月はコンバージョン率が高く、2025年08月はコンバージョン率が低い傾向にあります。",
    "",
    "・2026年04月のセッション数は7,457回、コンバージョン数は55件、コンバージョン率は0.74%でした。",
    "・2025年12月はセッション数5,818回に対し、コンバージョン数63件、コンバージョン率1.08%と、当期の中で最も高いコンバージョン率を記録しました。",
    "・2025年08月はセッション数5,486回に対し、コンバージョン数29件、コンバージョン率0.53%と、当期の中で最も低いコンバージョン率となりました。",
    "・エンゲージメント率（閲覧の活発さ）は概ね83%台後半で推移しており、全体的にユーザーの関心は高い状態が維持されています。",
])
EXTREME_AI = LONG_AI + "\n\n" + LONG_AI  # 倍長


def _summary(with_comp: bool = False) -> dict:
    return {
        "metrics": {"sessions": 7457, "totalUsers": 5205, "newUsers": 4803, "pageViews": 49766,
                    "engagementRate": 0.8337, "conversions": 55,
                    "clicks": 1856, "impressions": 96973, "ctr": 0.0191, "position": 12.123879842842854},
        "conversions": {},
    }


# ─── シナリオ定義 ─────────────────────────
def _conversions_data() -> dict:
    """conversion data (data: [{yearMonth, ev1, ev2, total}]) for builder."""
    months = ["2025-12", "2026-01", "2026-02", "2026-03", "2026-04"]
    return {
        "data": [
            {"yearMonth": m, "資料DL_Webサイト制作の流れ完了": 20 + i * 3,
             "新規お問い合わせ送信完了": 10 + i, "見積り依頼完了": 5 + i,
             "total": 35 + i * 5}
            for i, m in enumerate(months)
        ]
    }


def scenario_full(ai_text: str = LONG_AI) -> dict:
    """全シート + ユーザージャーニー + 流入KWファネル + 長文 AI"""
    return {
        "siteName": "テストサイト",
        "dateRange": {"from": "2026-04-01", "to": "2026-04-30"},
        "sheets": {key: _dynamic_sheet(5) for key in
                   ["monthly", "daily", "weekly", "hourly", "channels", "referrals",
                    "pages", "pageCategories", "landingPages", "fileDownloads", "externalLinks"]},
        "customSheets": {
            "siteUrl": "https://example.com",
            "summary": _summary(),
            "users": {"data": {"deviceCategory": [{"name": "desktop", "users": 100}]}},
            "conversions": _conversions_data(),
            "reverseFlows": [],
            "userJourney": _user_journey(),
            "keywordsFunnel": _keywords_funnel(),
            "improvements": [],
        },
        "aiAnalysis": {f"analysis/{key}": {"summary": ai_text}
                       for key in ["summary", "month", "day", "week", "hour", "users",
                                   "channels", "keywords", "referrals", "pages",
                                   "page-categories", "landing-pages", "file-downloads",
                                   "external-links", "conversions", "user-journey"]},
        "memos": {},
    }


def scenario_minimal() -> dict:
    """最小: summary + monthly のみ、ユーザージャーニー無し、AI 短文"""
    return {
        "siteName": "Minimum Site",
        "dateRange": {"from": "2026-03-01", "to": "2026-03-31"},
        "sheets": {"monthly": _dynamic_sheet(3)},
        "customSheets": {
            "siteUrl": "https://min.example.com",
            "summary": _summary(),
        },
        "aiAnalysis": {"analysis/summary": {"summary": SHORT_AI}},
        "memos": {},
    }


def scenario_extreme() -> dict:
    """極端: 比較期間 + 極長 AI + 大量データ"""
    payload = scenario_full(EXTREME_AI)
    payload["comparisonRange"] = {"from": "2026-03-01", "to": "2026-03-31"}
    payload["customSheets"]["compSummary"] = _summary()
    # detail_paths を 12 件に
    payload["customSheets"]["userJourney"] = _user_journey(num_paths=12)
    return payload


# ─── 検証ヘルパー ─────────────────────────
def verify_excel(path: Path) -> list[str]:
    """生成された .xlsx を読み戻し、構造の検証ログを返す。"""
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=False)
    logs = [f"\n=== Excel: {path.name} ==="]
    logs.append(f"  シート数: {len(wb.sheetnames)}")
    logs.append(f"  シート順: {wb.sheetnames}")

    # ユーザージャーニーシート詳細
    if "ユーザージャーニー" in wb.sheetnames:
        ws = wb["ユーザージャーニー"]
        logs.append(f"\n  [ユーザージャーニー] max_row={ws.max_row}, max_col={ws.max_column}")
        # 列幅
        widths = [ws.column_dimensions[chr(65 + i)].width for i in range(8)]
        logs.append(f"    列幅: {widths}")
        # セクションマーカー検出
        markers = []
        for r in range(1, ws.max_row + 1):
            v = ws.cell(r, 1).value
            if isinstance(v, str) and v.startswith("■"):
                markers.append(f"      R{r}: {v}")
        logs.append(f"    セクションマーカー:\n" + "\n".join(markers))
        # TOP3 の AI コメント行 (■ 主要ジャーニー TOP 3 の直後 5-7 行を確認)
        for r in range(1, min(ws.max_row + 1, 25)):
            v = ws.cell(r, 2).value
            if isinstance(v, str) and v.startswith("AI コメント"):
                logs.append(f"    R{r} AI コメント (B 列マージ): {v[:80]}...")
                logs.append(f"    R{r} 行高: {ws.row_dimensions[r].height}")

    # 流入キーワード元シート詳細
    if "流入キーワード元" in wb.sheetnames:
        ws = wb["流入キーワード元"]
        logs.append(f"\n  [流入キーワード元] max_row={ws.max_row}, max_col={ws.max_column}")
        widths = [ws.column_dimensions[chr(65 + i)].width for i in range(8)]
        logs.append(f"    列幅: {widths}")
        # クラスタ # 列の値 (■ 意味的クラスタ の直後にある #)
        cluster_marker_row = None
        for r in range(1, ws.max_row + 1):
            v = ws.cell(r, 1).value
            if isinstance(v, str) and "意味的クラスタ" in v:
                cluster_marker_row = r
                break
        if cluster_marker_row:
            # ヘッダーは marker+1, データは marker+2 から
            header_row = cluster_marker_row + 1
            logs.append(f"    クラスタヘッダー R{header_row}: A='{ws.cell(header_row, 1).value}', B='{ws.cell(header_row, 2).value}'")
            for r in range(cluster_marker_row + 2, min(cluster_marker_row + 8, ws.max_row + 1)):
                col_a = ws.cell(r, 1).value
                col_b = ws.cell(r, 2).value
                logs.append(f"    クラスタデータ R{r}: A='{col_a}', B='{col_b}'")

    return logs


def verify_pptx(path: Path) -> list[str]:
    """生成された .pptx を読み戻し、構造の検証ログを返す。"""
    from pptx import Presentation
    prs = Presentation(str(path))
    logs = [f"\n=== PPTX: {path.name} ==="]
    logs.append(f"  スライド数: {len(prs.slides)}")
    titles = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    title = txt[:60]
                    break
        titles.append(f"    {i:02d}. {title}")
    logs.extend(titles)
    return logs


# ─── 実行 ─────────────────────────
def run() -> None:
    scenarios = [
        ("scenario_1_full", scenario_full()),
        ("scenario_2_minimal", scenario_minimal()),
        ("scenario_3_extreme", scenario_extreme()),
    ]
    all_logs = []
    for name, payload in scenarios:
        # Excel
        xlsx_path = OUT_DIR / f"{name}.xlsx"
        with open(xlsx_path, "wb") as f:
            buf = io.BytesIO()
            build_excel_workbook(buf, payload)
            f.write(buf.getvalue())
        all_logs.extend(verify_excel(xlsx_path))

        # PPTX
        pptx_path = OUT_DIR / f"{name}.pptx"
        with open(pptx_path, "wb") as f:
            buf = io.BytesIO()
            build_pptx_presentation(buf, payload)
            f.write(buf.getvalue())
        all_logs.extend(verify_pptx(pptx_path))

    print("\n".join(all_logs))
    print(f"\n出力先: {OUT_DIR}")


if __name__ == "__main__":
    run()
