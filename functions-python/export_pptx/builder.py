"""
PPTX プレゼンテーション生成 — 旧 exportAnalysisToPptx.js のテンプレートを python-pptx に移植。

スライド構成:
  1. 表紙
  2. 全体サマリー
  3. トレンド分析（セクション区切り）
     - 月別（チャート / テーブル）
     - 日別（チャート / テーブル）
     - 曜日別
     - 時間帯別（チャート / テーブル）
  4. ユーザー分析（セクション区切り）
     - ドーナツ 4 本
     - 地域別ランキング
  5. 集客分析
     - チャネル / キーワード / 被リンク元
  6. コンテンツ分析
     - ページ別 / ページ分類別 / ランディング / ファイルDL / 外部リンク
  7. コンバージョン分析
     - 月次推移（チャート / テーブル）/ 逆算フロー
  8. Appendix（用語集）
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

from .helpers import (
    AI_MIN_H,
    CHANNEL_MAP,
    CHART_PALETTE,
    CHART_PALETTE_HEX,
    CONTENT_W,
    CONTENT_Y,
    Color,
    DAY_NAMES,
    FONT_FACE,
    FOOTER_Y,
    MARGIN_X,
    SLIDE_H,
    SLIDE_W,
    TITLE_H,
    TITLE_Y,
    calc_layout,
    calc_table_only_layout,
    clean_markdown,
    column_align,
    fmt_change,
    fmt_date,
    fmt_num,
    fmt_pct,
    fmt_year_month,
    format_cell_value,
    format_number,
    get_table_font_size,
    hex_to_rgb,
    set_cell_bg,
    set_cell_text,
)

from shared.metrics import (
    label_of,
    short_label_of,
    tooltip_of,
    comparison_label,
)


# ─── エントリポイント ────────────────────────────────────


class _Ctx:
    """スライド生成中の共有状態（ページ番号・日付ラベル等）。"""

    def __init__(self) -> None:
        self.slide_count = 0
        self.date_label = ""
        self.year_month = ""
        self.section_index = 0
        self.section_total = 0


def build_pptx_presentation(buffer: io.BytesIO, data: dict[str, Any]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    ctx = _Ctx()

    site_name = data.get("siteName") or "GrowReporter"
    date_range = data.get("dateRange") or {}
    comp_range = data.get("comparisonRange")
    sheets = data.get("sheets") or {}
    custom = data.get("customSheets") or {}
    ai = data.get("aiAnalysis") or {}
    memos = data.get("memos") or {}

    # 日付ラベル生成
    d_from = date_range.get("from") or ""
    d_to = date_range.get("to") or ""
    if d_from and d_to and len(d_to) >= 10:
        ctx.date_label = f"{d_from.replace('-', '.')}~{d_to[8:10]}"
    # セクション中表紙の右下に出す年月 (例: '26.03)
    ctx.year_month = ""
    if d_to and len(d_to) >= 7:
        ctx.year_month = f"'{d_to[2:4]}.{d_to[5:7]}"
    # セクション番号管理 (中表紙の通し番号 + 総数)
    # 各セクションは中身が 1 件以上ある場合のみ中表紙を出すため動的計算
    ctx.section_index = 0
    def _has_rows(key: str) -> bool:
        s = sheets.get(key)
        return bool(s and s.get("rows"))

    has_analysis = bool(custom.get("summary") or custom.get("users"))
    has_timeline = any(_has_rows(k) for k in ("monthly", "daily", "weekly", "hourly"))
    kw_v2 = custom.get("keywordsFunnel")
    has_marketing = bool(
        _has_rows("channels")
        or (kw_v2 and (kw_v2.get("funnel") or kw_v2.get("keywords")))
        or _has_rows("referrals")
    )
    user_journey = custom.get("userJourney")
    has_pages = bool(
        any(_has_rows(k) for k in ("pages", "pageCategories", "landingPages", "fileDownloads", "externalLinks"))
        or (user_journey and user_journey.get("nodes"))
    )
    has_conversion = bool(
        (custom.get("conversions") and custom["conversions"].get("data"))
        or custom.get("reverseFlows")
    )
    has_improvements = bool(custom.get("improvements"))
    ctx.section_total = sum([
        int(has_analysis), int(has_timeline), int(has_marketing), int(has_pages),
        int(has_conversion), int(has_improvements), 1,  # Appendix は常に出す
    ])

    # ナビ順 (アプリと完全一致):
    #  分析: 全体サマリー → ユーザー属性
    #  時系列: 月別 → 日別 → 曜日別 → 時間帯別
    #  集客: 集客チャネル → 流入キーワード元 → 被リンク元
    #  ページ: ページ別 → ページ分類別 → ランディングページ →
    #         ファイルダウンロード → 外部リンククリック → ユーザージャーニー
    #  コンバージョン: コンバージョン一覧 → 逆算フロー
    #  (改善アクション)
    #  Appendix

    # 1. 表紙
    _create_cover_slide(prs, ctx, site_name, custom.get("siteUrl") or "", date_range, comp_range)

    # ── 分析 ──
    if has_analysis:
        _create_section_divider(prs, ctx, "分析", eyebrow="ANALYSIS")
    _create_summary_slide(
        prs, ctx,
        summary=custom.get("summary"),
        comp_summary=custom.get("compSummary"),
        kpi_settings=custom.get("kpiSettings"),
        monthly_delta=custom.get("monthlyDelta"),
        ai_data=ai.get("analysis/summary"),
        memos=memos.get("analysis/summary") or memos.get("summary"),
    )
    _create_users_donut_slide(prs, ctx, demographics=custom.get("users"))
    _create_users_region_slide(
        prs, ctx,
        demographics=custom.get("users"),
        ai_data=ai.get("analysis/users"),
        memos=memos.get("analysis/users") or memos.get("users"),
    )

    # ── 時系列 ──
    if has_timeline:
        _create_section_divider(prs, ctx, "時系列", eyebrow="TIMELINE")
    _create_monthly_slides(
        prs, ctx,
        monthly_sheet=sheets.get("monthly"),
        ai_data=ai.get("analysis/month"),
        memos=memos.get("analysis/month"),
    )
    _create_daily_slides(
        prs, ctx,
        daily_sheet=sheets.get("daily"),
        ai_data=ai.get("analysis/day"),
        memos=memos.get("analysis/day") or memos.get("day"),
    )
    _create_weekly_slide(
        prs, ctx,
        weekly_sheet=sheets.get("weekly"),
        ai_data=ai.get("analysis/week"),
        memos=memos.get("analysis/week") or memos.get("week"),
    )
    _create_hourly_slides(
        prs, ctx,
        hourly_sheet=sheets.get("hourly"),
        ai_data=ai.get("analysis/hour"),
        memos=memos.get("analysis/hour") or memos.get("hour"),
    )

    # ── 集客 ──
    if has_marketing:
        _create_section_divider(prs, ctx, "集客", eyebrow="MARKETING")
    _create_channels_slide(
        prs, ctx,
        channels_sheet=sheets.get("channels"),
        ai_data=ai.get("analysis/channels"),
        memos=memos.get("analysis/channels") or memos.get("channels"),
    )
    # 流入キーワード元 = V2 ファネル (旧 Top 20 はナビから外れたため非表示)
    _create_keywords_funnel_slide(
        prs, ctx,
        keywords_v2=custom.get("keywordsFunnel"),
        ai_data=ai.get("analysis/keywords"),
        memos=memos.get("analysis/keywords") or memos.get("keywords"),
    )
    _create_referrals_slide(
        prs, ctx,
        referrals_sheet=sheets.get("referrals"),
        ai_data=ai.get("analysis/referrals"),
        memos=memos.get("analysis/referrals") or memos.get("referrals"),
    )

    # ── ページ ──
    if has_pages:
        _create_section_divider(prs, ctx, "ページ", eyebrow="PAGES")
    _create_pages_slide(
        prs, ctx,
        pages_sheet=sheets.get("pages"),
        ai_data=ai.get("analysis/pages"),
        memos=memos.get("analysis/pages") or memos.get("pages"),
    )
    _create_page_categories_slide(
        prs, ctx,
        page_categories_sheet=sheets.get("pageCategories"),
        ai_data=ai.get("analysis/page-categories"),
        memos=memos.get("analysis/page-categories") or memos.get("pageCategories"),
    )
    _create_landing_pages_slide(
        prs, ctx,
        landing_sheet=sheets.get("landingPages"),
        ai_data=ai.get("analysis/landing-pages"),
        memos=memos.get("analysis/landing-pages") or memos.get("landingPages"),
    )
    _create_file_downloads_slide(
        prs, ctx,
        sheet=sheets.get("fileDownloads"),
        ai_data=ai.get("analysis/file-downloads"),
        memos=memos.get("analysis/file-downloads") or memos.get("fileDownloads"),
    )
    _create_external_links_slide(
        prs, ctx,
        sheet=sheets.get("externalLinks"),
        ai_data=ai.get("analysis/external-links"),
        memos=memos.get("analysis/external-links") or memos.get("externalLinks"),
    )
    _create_user_journey_slide(
        prs, ctx,
        user_journey=custom.get("userJourney"),
        ai_data=ai.get("analysis/user-journey"),
        memos=memos.get("analysis/user-journey") or memos.get("userJourney"),
    )

    # ── コンバージョン ──
    if has_conversion:
        _create_section_divider(prs, ctx, "コンバージョン", eyebrow="CONVERSION")
    _create_conversions_slides(
        prs, ctx,
        conversions=custom.get("conversions"),
        ai_data=ai.get("analysis/conversions"),
        memos=memos.get("analysis/conversions") or memos.get("conversions"),
    )
    _create_reverse_flow_slide(
        prs, ctx,
        reverse_flows=custom.get("reverseFlows"),
        ai_data=ai.get("analysis/reverse-flow"),
        memos=memos.get("analysis/reverse-flow") or memos.get("reverseFlow"),
    )

    # ── 改善アクション (improvements がある場合のみ) ──
    if custom.get("improvements"):
        _create_section_divider(prs, ctx, "改善アクション", eyebrow="ACTION")
        _create_improvements_slide(prs, ctx, custom.get("improvements"))

    # ── Appendix (用語解説) ──
    _create_section_divider(prs, ctx, "用語解説", eyebrow="APPENDIX")
    _create_appendix_slide(prs, ctx)

    prs.save(buffer)


# ─── 共通フレーム ────────────────────────────────────────


def _add_slide_title(slide, ctx: _Ctx, title: str) -> None:
    tb = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(TITLE_Y),
        Inches(CONTENT_W - 2.0), Inches(TITLE_H),
    )
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_FACE
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = Color.DARK
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # タイトル下の横線（1pt 相当、PowerPoint が扱える最小幅に合わせる）
    line = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(MARGIN_X), Inches(TITLE_Y + TITLE_H + 0.03),
        Inches(MARGIN_X + CONTENT_W), Inches(TITLE_Y + TITLE_H + 0.03),
    )
    line.line.color.rgb = Color.DARK
    line.line.width = Pt(1)
    _disable_shape_shadow(line)

    # 日付ラベル（右上）
    if ctx.date_label:
        dl = slide.shapes.add_textbox(
            Inches(SLIDE_W - MARGIN_X - 1.6), Inches(TITLE_Y + 0.08),
            Inches(1.6), Inches(0.28),
        )
        dtf = dl.text_frame
        dtf.margin_left = 0
        dtf.margin_right = 0
        dp = dtf.paragraphs[0]
        dp.alignment = PP_ALIGN.RIGHT
        drun = dp.add_run()
        drun.text = ctx.date_label
        drun.font.name = FONT_FACE
        drun.font.size = Pt(9)
        drun.font.color.rgb = Color.PRIMARY
        dl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_slide_footer(slide, ctx: _Ctx) -> None:
    """全スライド共通フッター (表紙除く):
    左に © グローレポータ / 右に P.NN。スライドコンテンツ幅 (MARGIN_X ~ SLIDE_W - MARGIN_X) に揃える。
    """
    ctx.slide_count += 1
    foot_y = SLIDE_H - 0.32
    foot_h = 0.28
    page_w = 0.8

    # 左: © (コンテンツ左端 = MARGIN_X)
    brand_tb = slide.shapes.add_textbox(
        Inches(MARGIN_X), Inches(foot_y),
        Inches(CONTENT_W - page_w - 0.2), Inches(foot_h),
    )
    btf = brand_tb.text_frame
    btf.margin_left = 0
    btf.margin_right = 0
    btf.margin_top = 0
    btf.margin_bottom = 0
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = "© グローレポータ All rights reserved."
    br.font.name = FONT_FACE
    br.font.size = Pt(9)
    br.font.color.rgb = Color.SUB_TEXT

    # 右: P.NN (コンテンツ右端 = SLIDE_W - MARGIN_X)
    page_tb = slide.shapes.add_textbox(
        Inches(SLIDE_W - MARGIN_X - page_w), Inches(foot_y),
        Inches(page_w), Inches(foot_h),
    )
    ptf = page_tb.text_frame
    ptf.margin_left = 0
    ptf.margin_right = 0
    ptf.margin_top = 0
    ptf.margin_bottom = 0
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = ptf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    pr = pp.add_run()
    pr.text = f"P.{ctx.slide_count:02d}"
    pr.font.name = FONT_FACE
    pr.font.size = Pt(9)
    pr.font.color.rgb = Color.SUB_TEXT


def _add_section_label(slide, ctx: _Ctx) -> None:
    """セクション区切りスライドは footer 色を白で描画。"""
    ctx.slide_count += 1
    tb = slide.shapes.add_textbox(
        Inches(SLIDE_W - 1.5), Inches(FOOTER_Y),
        Inches(1.2), Inches(0.3),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(ctx.slide_count)
    run.font.name = FONT_FACE
    run.font.size = Pt(8)
    run.font.color.rgb = Color.WHITE
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


# ─── 表（テーブル）ヘルパー ─────────────────────────────


def _build_table(
    slide,
    headers: list[dict],
    rows: list[list[Any]],
    col_widths: list[float],
    x: float,
    y: float,
    h: float,
    *,
    font_size: float | None = None,
    header_align_default: str = "center",
):
    """旧 buildTable 相当。headers = [{label, align}], rows は文字列 / 数値。"""
    total_w = sum(col_widths)
    cols = len(headers)
    rows_n = len(rows) + 1

    shape = slide.shapes.add_table(
        rows_n, cols,
        Inches(x), Inches(y),
        Inches(total_w), Inches(h),
    )
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    table.last_row = False
    table.first_col = False
    table.last_col = False

    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    row_h = h / rows_n
    for i in range(rows_n):
        table.rows[i].height = Inches(row_h)

    if font_size is None:
        font_size = get_table_font_size(len(rows))

    # ヘッダー（見出し行は常にセンター揃え）
    for c, h_def in enumerate(headers):
        cell = table.cell(0, c)
        set_cell_bg(cell, Color.PRIMARY)
        set_cell_text(
            cell, h_def.get("label", ""),
            font_size=max(font_size, 8.5), bold=True,
            color=Color.WHITE,
            align="center",
            vertical="middle",
        )
        _set_cell_border(cell, Color.BORDER, width_pt=0.5)

    # データ行
    for r_idx, row in enumerate(rows):
        bg = Color.ALT_ROW if r_idx % 2 == 0 else Color.WHITE
        for c, val in enumerate(row):
            cell = table.cell(r_idx + 1, c)
            set_cell_bg(cell, bg)
            align = headers[c].get("align") if c < len(headers) else None
            if not align:
                align = "right" if isinstance(val, (int, float)) else "left"
            set_cell_text(
                cell, _cell_str(val),
                font_size=font_size,
                color=Color.DARK,
                align=align,
                vertical="middle",
            )
            _set_cell_border(cell, Color.BORDER, width_pt=0.5)

    return table


def _cell_str(val: Any) -> str:
    if val is None or val == "":
        return "-"
    if isinstance(val, float):
        if abs(val - int(val)) < 1e-9:
            return f"{int(val):,}"
        return f"{val:,.2f}"
    if isinstance(val, int):
        return f"{val:,}"
    return str(val)


def _disable_shape_shadow(shape) -> None:
    """図形のテーマ由来のエフェクト（影）を無効化する。

    <p:style><a:effectRef/></p:style> で参照されるテーマ効果を <a:effectLst/> で上書きする。
    <a:effectLst/> は <p:spPr> の EG_EffectProperties 位置（ln の後、scene3d/sp3d/extLst の前）に挿入する。
    """
    sp_pr = shape._element.spPr
    # 既存の effectLst/effectDag を削除
    for tag in ("a:effectLst", "a:effectDag"):
        for e in sp_pr.findall(qn(tag)):
            sp_pr.remove(e)
    # 挿入位置: scene3d/sp3d/extLst の直前（無ければ末尾）
    effect_lst = etree.Element(qn("a:effectLst"))
    successors = ("a:scene3d", "a:sp3d", "a:extLst")
    insert_idx = len(sp_pr)
    for i, child in enumerate(sp_pr):
        if etree.QName(child).localname in [qn(t).split("}")[-1] for t in successors]:
            insert_idx = i
            break
    sp_pr.insert(insert_idx, effect_lst)


def _set_cell_border(cell, color: "RGBColor", width_pt: float = 0.5) -> None:
    """4 方向のセル罫線を設定。

    OOXML スキーマ (CT_TableCellProperties) は lnL → lnR → lnT → lnB → 塗り の順を要求するため、
    tcPr の先頭に挿入する（塗りが既に追加されていても正しい順序を保つ）。
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    hex_color = "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])
    width_emu = int(width_pt * 12700)

    tag_order = ("a:lnL", "a:lnR", "a:lnT", "a:lnB")

    # 既存の罫線要素を削除
    for tag in tag_order:
        for e in tcPr.findall(qn(tag)):
            tcPr.remove(e)

    # L → R → T → B の順で先頭に挿入
    for i, tag in enumerate(tag_order):
        ln = etree.Element(qn(tag))
        ln.set("w", str(width_emu))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", hex_color)
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "solid")
        etree.SubElement(ln, qn("a:round"))
        tcPr.insert(i, ln)


# ─── AI 分析 + メモ フッター ─────────────────────────────


def _add_ai_and_memo_footer(
    slide,
    ai_data: dict | None,
    memos: list | None,
    y: float,
    max_h: float,
) -> None:
    if not ai_data or not ai_data.get("summary"):
        return
    ai_text = clean_markdown(ai_data.get("summary"))
    if not ai_text:
        return

    # AI 分析ラベル（紫）— PowerPoint の描画互換性のため通常の矩形を使用
    label = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN_X), Inches(y),
        Inches(1.2), Inches(0.28),
    )
    label.fill.solid()
    label.fill.fore_color.rgb = Color.ACCENT
    label.line.fill.background()
    _disable_shape_shadow(label)
    ltf = label.text_frame
    ltf.margin_left = Emu(36576)
    ltf.margin_right = Emu(36576)
    ltf.margin_top = 0
    ltf.margin_bottom = 0
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lrun = lp.add_run()
    lrun.text = "AI分析"
    lrun.font.name = FONT_FACE
    lrun.font.size = Pt(10)
    lrun.font.bold = True
    lrun.font.color.rgb = Color.WHITE
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # AI 本文
    text_y = y + 0.35
    text_h = max(max_h - 0.35, 0.3)
    tb = slide.shapes.add_textbox(
        Inches(MARGIN_X + 0.1), Inches(text_y),
        Inches(CONTENT_W - 0.2), Inches(text_h),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(18288)
    tf.margin_right = Emu(18288)

    # テキストを段落ごとに追加（行間 1.5 倍）
    paragraphs = ai_text.split("\n")
    for idx, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.5
        run = p.add_run()
        run.text = line
        run.font.name = FONT_FACE
        run.font.size = Pt(10)
        run.font.color.rgb = Color.DARK


def _create_ai_slide(
    prs, ctx: "_Ctx", parent_title: str,
    ai_data: dict | None, memos: list | None,
) -> None:
    """親スライドの直後に AI 分析専用スライドを追加する。
    AI データもメモも無い場合はスキップ。
    """
    has_ai = bool(ai_data and ai_data.get("summary"))
    has_memo = bool(memos and len(memos) > 0)
    if not has_ai and not has_memo:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, f"{parent_title} — AI分析")

    # コンテンツ全域を AI に使用
    ai_y = CONTENT_Y
    ai_h = FOOTER_Y - CONTENT_Y - 0.1
    _add_ai_and_memo_footer(slide, ai_data, memos, ai_y, ai_h)


# ─── 1. 表紙 ────────────────────────────────────────────


def _create_cover_slide(
    prs: Presentation, ctx: _Ctx,
    site_name: str, site_url: str,
    date_range: dict | None, comp_range: dict | None,
) -> None:
    """表紙: 上部 (REPORT 番号 + ロゴ) / 中央 (大タイトル) / 右側 (縦書き) / 下部 (4 列情報)"""
    from datetime import datetime as _dt, timedelta, timezone

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.BG_SOFT)
    # 表紙はページ番号管理から除外 (slide_count を増やさない)

    d_from = (date_range or {}).get("from", "") or ""
    d_to = (date_range or {}).get("to", "") or ""
    end_ym = ""
    if d_to and len(d_to) >= 7:
        end_ym = d_to[:7].replace("-", ".")  # 2026.03

    # 期間文字列 (同じ年月なら短縮形)
    period_str = ""
    if d_from and d_to:
        if d_from[:7] == d_to[:7] and len(d_to) >= 10:
            period_str = f"{d_from[:7].replace('-', '.')}.{d_from[8:10]}—{d_to[8:10]}"
        else:
            period_str = f"{d_from.replace('-', '.')}—{d_to.replace('-', '.')}"

    JST = timezone(timedelta(hours=9))
    created_str = _dt.now(JST).strftime("%Y/%m/%d")

    left_x = 0.7
    right_x = SLIDE_W - 0.7

    # ─── 上部左: ─ REPORT / YYYY.MM ──
    rep_y = 0.5
    rep_h = 0.30
    rep_tb = slide.shapes.add_textbox(
        Inches(left_x + 0.4), Inches(rep_y),
        Inches(5.0), Inches(rep_h),
    )
    rtf = rep_tb.text_frame
    rtf.margin_left = 0
    rtf.margin_right = 0
    rtf.margin_top = 0
    rtf.margin_bottom = 0
    rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rp = rtf.paragraphs[0]
    rr = rp.add_run()
    rr.text = f"REPORT / {end_ym}"
    rr.font.name = FONT_FACE
    rr.font.size = Pt(10)
    rr.font.color.rgb = Color.PRIMARY

    # 装飾線 (テキスト中央高さに揃える)
    deco_y = rep_y + rep_h / 2
    deco = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left_x), Inches(deco_y),
        Inches(left_x + 0.25), Inches(deco_y),
    )
    deco.line.color.rgb = Color.PRIMARY
    deco.line.width = Pt(1.2)
    _disable_shape_shadow(deco)

    # ─── ロゴ画像 (タグラインはロゴ内にあるため別表示なし) ──
    import os as _os
    _logo_path = _os.path.join(
        _os.path.dirname(_os.path.dirname(_os.path.abspath(__file__))),
        "assets", "logo.png",
    )
    if _os.path.exists(_logo_path):
        logo_w_in = 2.3
        logo_h_in = logo_w_in * (383 / 1600)
        slide.shapes.add_picture(
            _logo_path,
            Inches(left_x), Inches(1.05),
            width=Inches(logo_w_in), height=Inches(logo_h_in),
        )

    # ─── 中央: 大タイトル「分析レポート」 + サイト名 ──
    title_y = 2.6
    title_tb = slide.shapes.add_textbox(
        Inches(left_x), Inches(title_y),
        Inches(SLIDE_W - left_x - 1.2), Inches(1.4),
    )
    ttf = title_tb.text_frame
    ttf.margin_left = 0
    ttf.margin_right = 0
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    tr.text = "分析レポート"
    tr.font.name = FONT_FACE
    tr.font.size = Pt(64)
    tr.font.bold = True
    tr.font.color.rgb = Color.DARK

    # サイト名 (タイトル直下、タイトルに寄せて配置)
    if site_name:
        sub_tb = slide.shapes.add_textbox(
            Inches(left_x), Inches(title_y + 1.25),
            Inches(SLIDE_W - left_x - 1.2), Inches(0.4),
        )
        stf = sub_tb.text_frame
        stf.margin_left = 0
        stf.margin_right = 0
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.add_run()
        sr.text = site_name
        sr.font.name = FONT_FACE
        sr.font.size = Pt(14)
        sr.font.color.rgb = Color.SUB_TEXT

    # ─── 下部: 横線 + 4 カラム情報 (フッターとの距離を確保) ──
    bottom_y = SLIDE_H - 1.8
    bline = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left_x), Inches(bottom_y),
        Inches(right_x), Inches(bottom_y),
    )
    bline.line.color.rgb = Color.BORDER
    bline.line.width = Pt(0.5)
    _disable_shape_shadow(bline)

    info_cols = [
        ("01・SITE", site_name or ""),
        ("02・URL", site_url or ""),
        ("03・PERIOD", period_str or ""),
        ("04・CREATED", created_str),
    ]
    col_w = (right_x - left_x) / 4
    for i, (label, value) in enumerate(info_cols):
        x = left_x + col_w * i
        # ラベル
        ltb = slide.shapes.add_textbox(
            Inches(x), Inches(bottom_y + 0.18),
            Inches(col_w - 0.2), Inches(0.28),
        )
        ltf = ltb.text_frame
        ltf.margin_left = 0
        ltf.margin_right = 0
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        lr.font.name = FONT_FACE
        lr.font.size = Pt(8)
        lr.font.color.rgb = Color.SUB_TEXT
        # 値
        vtb = slide.shapes.add_textbox(
            Inches(x), Inches(bottom_y + 0.5),
            Inches(col_w - 0.2), Inches(0.4),
        )
        vtf2 = vtb.text_frame
        vtf2.margin_left = 0
        vtf2.margin_right = 0
        vtf2.word_wrap = True
        vp2 = vtf2.paragraphs[0]
        vr2 = vp2.add_run()
        vr2.text = value
        vr2.font.name = FONT_FACE
        vr2.font.size = Pt(11)
        vr2.font.color.rgb = Color.DARK

    # 共通フッター (© + P.01) — 表紙も含めて通し番号
    _add_slide_footer(slide, ctx)


def _fill_slide_bg(slide, color) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ─── 2. 全体サマリー ─────────────────────────────────────


def _create_summary_slide(
    prs: Presentation, ctx: _Ctx,
    summary: dict | None,
    comp_summary: dict | None,
    kpi_settings: dict | None,
    ai_data: dict | None,
    memos: list | None,
    monthly_delta: dict | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "全体サマリー")

    if not summary or not summary.get("metrics"):
        _add_slide_footer(slide, ctx)
        return

    m = summary["metrics"]
    cm = (comp_summary or {}).get("metrics") or {}
    has_comp = bool(cm)

    # 項目: 現在の 6 コア + 4 GSC を維持
    # 上段 (大カード) = セッション / コンバージョン / エンゲージメント率 の 3 主要指標
    # 下段 (小カード) = ユーザー / 新規 / PV / クリック / 表示回数 / CTR / 掲載順位
    # mom_key は monthlyDelta の lookup キー (None なら前月比表示なし)
    top_items = [
        (short_label_of("sessions"), m.get("sessions"), cm.get("sessions"), "number", "sessions"),
        (short_label_of("conversions"), m.get("conversions"), cm.get("conversions"), "number", "conversions"),
        (short_label_of("engagementRate"), m.get("engagementRate"), cm.get("engagementRate"), "percent", "engagementRate"),
    ]
    bottom_items = [
        (short_label_of("totalUsers"), m.get("totalUsers"), cm.get("totalUsers"), "number", "totalUsers"),
        (short_label_of("newUsers"), m.get("newUsers"), cm.get("newUsers"), "number", "newUsers"),
        (short_label_of("screenPageViews"), m.get("pageViews"), cm.get("pageViews"), "number", "pageViews"),
    ]
    if m.get("clicks") is not None or m.get("impressions") is not None:
        bottom_items.extend([
            (short_label_of("clicks"), m.get("clicks"), cm.get("clicks"), "number", None),
            (short_label_of("impressions"), m.get("impressions"), cm.get("impressions"), "number", None),
            (short_label_of("ctr"), m.get("ctr"), cm.get("ctr"), "percent", None),
            (short_label_of("position"), m.get("position"), cm.get("position"), "decimal", None),
        ])

    def _fmt_val(val, kind):
        if val is None or val == "":
            return "-"
        if kind == "percent":
            try:
                v = float(val)
                v = v * 100 if abs(v) <= 1 else v
                return f"{v:.1f}"
            except (ValueError, TypeError):
                return "-"
        if kind == "decimal":
            try:
                return f"{float(val):.1f}"
            except (ValueError, TypeError):
                return "-"
        return format_number(val)

    def _fmt_delta(cur, prev, kind):
        """変化率 (差分) を返す。None なら表示しない。"""
        if cur is None or prev is None:
            return None
        try:
            c, p = float(cur), float(prev)
        except (ValueError, TypeError):
            return None
        if kind == "percent":
            # 割合系はパーセントポイント差
            cv = c * 100 if abs(c) <= 1 else c
            pv = p * 100 if abs(p) <= 1 else p
            diff = cv - pv
            sign = "+" if diff >= 0 else ""
            return (f"{sign}{diff:.1f}%", diff >= 0)
        # 数値系は前期比
        if p == 0:
            return None
        diff = (c - p) / p * 100
        sign = "+" if diff >= 0 else ""
        return (f"{sign}{diff:.1f}%", diff >= 0)

    POS_COLOR = RGBColor(0x10, 0xB9, 0x81)  # green
    NEG_COLOR = RGBColor(0xEF, 0x44, 0x44)  # red
    CARD_BG = RGBColor(0xF3, 0xF4, 0xF6)    # very light gray

    def _unit_suffix(kind):
        return "%" if kind == "percent" else ""

    def _draw_card(x, y, w, h, label, value_str, unit, delta=None, mom=None, big=True):
        """1 枚のカード (rounded rectangle) を描画。delta=比較期間変化率 / mom=前月比。"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(w), Inches(h),
        )
        # 控えめな角丸 (default ~0.17 → 0.05 で小さめ角丸)
        try:
            card.adjustments[0] = 0.05
        except Exception:
            pass
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()
        _disable_shape_shadow(card)

        # ラベル
        label_tb = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.15),
            Inches(w - 0.4), Inches(0.28),
        )
        ltf = label_tb.text_frame
        ltf.margin_left = 0
        ltf.margin_right = 0
        ltf.margin_top = 0
        ltf.margin_bottom = 0
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        lr.font.name = FONT_FACE
        lr.font.size = Pt(10)
        lr.font.color.rgb = Color.SUB_TEXT

        # 値 (数値 + 単位)
        val_y = y + 0.45
        val_tb = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(val_y),
            Inches(w - 0.4), Inches(0.7 if big else 0.55),
        )
        vtf = val_tb.text_frame
        vtf.margin_left = 0
        vtf.margin_right = 0
        vtf.margin_top = 0
        vtf.margin_bottom = 0
        vp = vtf.paragraphs[0]
        vr = vp.add_run()
        vr.text = value_str
        vr.font.name = FONT_FACE
        vr.font.size = Pt(32 if big else 22)
        vr.font.bold = True
        vr.font.color.rgb = Color.DARK
        if unit:
            ur = vp.add_run()
            ur.text = unit
            ur.font.name = FONT_FACE
            ur.font.size = Pt(16 if big else 12)
            ur.font.color.rgb = Color.DARK

        # フッターに 比較期間差分 と 前月比 を併記 (利用可能なもののみ)
        footer_items = []
        if delta is not None:
            delta_str, is_positive = delta
            footer_items.append(("前期比", delta_str, is_positive))
        if mom is not None:
            mom_str, is_positive = mom
            footer_items.append(("前月比", mom_str, is_positive))

        if footer_items:
            foot_tb = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + h - 0.32),
                Inches(w - 0.4), Inches(0.25),
            )
            ftf = foot_tb.text_frame
            ftf.margin_left = 0
            ftf.margin_right = 0
            ftf.margin_top = 0
            ftf.margin_bottom = 0
            fp = ftf.paragraphs[0]
            for idx, (prefix, text, is_pos) in enumerate(footer_items):
                if idx > 0:
                    sep = fp.add_run()
                    sep.text = "  "
                    sep.font.name = FONT_FACE
                    sep.font.size = Pt(8)
                # ラベル (グレー)
                lbl = fp.add_run()
                lbl.text = f"{prefix} "
                lbl.font.name = FONT_FACE
                lbl.font.size = Pt(8)
                lbl.font.color.rgb = Color.SUB_TEXT
                # 値 (緑/赤)
                val = fp.add_run()
                val.text = text
                val.font.name = FONT_FACE
                val.font.size = Pt(9)
                val.font.bold = True
                val.font.color.rgb = POS_COLOR if is_pos else NEG_COLOR

    md = monthly_delta or {}

    def _mom_pair(key):
        """monthlyDelta から前月比を取り出し整形。"""
        if not key or key not in md:
            return None
        mom_v = (md.get(key) or {}).get("mom")
        if mom_v is None:
            return None
        try:
            pct = float(mom_v) * 100
        except (ValueError, TypeError):
            return None
        sign = "+" if pct >= 0 else ""
        return (f"{sign}{pct:.1f}%", pct >= 0)

    # ─── 上段: 大カード 3 枚 ──
    top_y = CONTENT_Y
    gap = 0.15
    n_top = len(top_items)
    top_card_w = (CONTENT_W - gap * (n_top - 1)) / n_top
    top_card_h = 1.5
    for i, (label, cur, prev, kind, mk) in enumerate(top_items):
        x = MARGIN_X + (top_card_w + gap) * i
        delta = _fmt_delta(cur, prev, kind) if has_comp else None
        mom = _mom_pair(mk)
        _draw_card(x, top_y, top_card_w, top_card_h, label, _fmt_val(cur, kind), _unit_suffix(kind), delta=delta, mom=mom, big=True)

    # ─── 下段: 小カード ──
    bot_y = top_y + top_card_h + 0.2
    n_bot = len(bottom_items)
    bot_card_w = (CONTENT_W - gap * (n_bot - 1)) / n_bot
    bot_card_h = 1.2
    for i, (label, cur, prev, kind, mk) in enumerate(bottom_items):
        x = MARGIN_X + (bot_card_w + gap) * i
        mom = _mom_pair(mk)
        _draw_card(x, bot_y, bot_card_w, bot_card_h, label, _fmt_val(cur, kind), _unit_suffix(kind), delta=None, mom=mom, big=False)

    kpi_end_y = bot_y + bot_card_h + 0.2

    # 目標達成状況
    if kpi_settings and kpi_settings.get("kpiList"):
        active_kpis = [k for k in kpi_settings["kpiList"] if k.get("isActive")]
        if active_kpis:
            tb = slide.shapes.add_textbox(
                Inches(MARGIN_X), Inches(kpi_end_y),
                Inches(CONTENT_W), Inches(0.3),
            )
            tf = tb.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "目標達成状況"
            run.font.name = FONT_FACE
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = Color.DARK
            kpi_end_y += 0.35

            headers = [
                {"label": "指標", "align": "left"},
                {"label": "目標値", "align": "right"},
                {"label": "実績値", "align": "right"},
                {"label": "達成率", "align": "right"},
            ]
            rows_out = []
            for k in active_kpis:
                actual = _get_kpi_actual(k, summary)
                tgt = k.get("target")
                try:
                    rate = f"{(float(actual) / float(tgt)) * 100:.1f}%" if tgt and actual is not None else "-"
                except (ValueError, TypeError, ZeroDivisionError):
                    rate = "-"
                rows_out.append([
                    k.get("label") or k.get("metric") or "",
                    format_number(tgt),
                    format_number(actual),
                    rate,
                ])
            tbl_h = 0.36 * (len(rows_out) + 1)  # 目標表も少し低く（旧 0.4）
            _build_table(
                slide, headers, rows_out, [3.5, 2.2, 2.2, 2.2],
                x=MARGIN_X, y=kpi_end_y, h=tbl_h,
                font_size=10,
            )
            kpi_end_y += tbl_h + 0.1

    _add_slide_footer(slide, ctx)

    # AI 分析は独立スライドへ (本スライドの下にはみ出す問題を回避)
    _create_ai_slide(prs, ctx, "全体サマリー", ai_data, memos)


def _get_kpi_actual(kpi: dict, summary: dict | None) -> Any:
    if not summary or not summary.get("metrics"):
        return None
    m = summary["metrics"]
    mapping = {
        "sessions": m.get("sessions"),
        "totalUsers": m.get("totalUsers"),
        "newUsers": m.get("newUsers"),
        "screenPageViews": m.get("pageViews"),
        "engagementRate": (m.get("engagementRate") * 100) if isinstance(m.get("engagementRate"), (int, float)) else None,
        "conversions": m.get("conversions"),
        "clicks": m.get("clicks"),
        "impressions": m.get("impressions"),
    }
    if kpi.get("isConversion") and kpi.get("eventName"):
        cv = summary.get("conversions") or {}
        return cv.get(kpi["eventName"], 0)
    return mapping.get(kpi.get("metric"))


# ─── セクション区切り ────────────────────────────────────


SECTION_DESCRIPTIONS = {
    "トレンド分析": "月別・日別・曜日別・時間帯別の切り口で、訪問・成果の推移を可視化します。\nどの期間にどのくらいの動きがあったかを把握できます。",
    "ユーザー分析": "デバイス・年齢・性別・地域別のユーザー属性とアクセス傾向をまとめています。",
    "集客分析": "チャネル・キーワード・参照元から、流入の量と質を多角的に分析します。",
    "コンテンツ分析": "ページ別・分類別・LP・ファイル DL・外部リンクの利用状況を確認できます。",
    "コンバージョン分析": "コンバージョン推移と、フォーム到達までの逆算フローによるファネル分析を確認できます。",
    "改善アクション": "AI が抽出した改善提案を、優先度・カテゴリ・期待効果と合わせて一覧化します。\nモックアップ生成済みの提案は、共有 URL から After 画面をその場で確認できます。",
    "用語解説": "用語・指標の説明と補足情報をまとめています。",
}


def _create_section_divider(prs: Presentation, ctx: _Ctx, title: str, *, eyebrow: str = "基本レポート") -> None:
    """中表紙: 薄背景 + 巨大薄色番号 + SECTION/CHAPTER ラベル + タイトル + 説明 + 専用フッター"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.BG_SOFT)
    ctx.section_index += 1
    sec_idx = ctx.section_index
    sec_total = ctx.section_total or 6
    sec_num_str = f"{sec_idx:02d}"
    total_num_str = f"{sec_total:02d}"

    left_x = 0.7
    right_x = SLIDE_W - 0.7

    # ─── 上部左: ─ SECTION / NN / 基本レポート ──
    top_y = 0.45
    top_h = 0.30
    sec_tb = slide.shapes.add_textbox(
        Inches(left_x + 0.4), Inches(top_y),
        Inches(5.0), Inches(top_h),
    )
    stf = sec_tb.text_frame
    stf.margin_left = 0
    stf.margin_right = 0
    stf.margin_top = 0
    stf.margin_bottom = 0
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = f"SECTION / {sec_num_str} / {eyebrow}"
    sr.font.name = FONT_FACE
    sr.font.size = Pt(10)
    sr.font.color.rgb = Color.PRIMARY

    # 装飾線 (テキスト中央高さに揃える)
    deco_y = top_y + top_h / 2
    deco = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left_x), Inches(deco_y),
        Inches(left_x + 0.25), Inches(deco_y),
    )
    deco.line.color.rgb = Color.PRIMARY
    deco.line.width = Pt(1.2)
    _disable_shape_shadow(deco)

    # ─── 上部右: CHAPTER NN / NN ──
    chap_tb = slide.shapes.add_textbox(
        Inches(SLIDE_W - 3.0), Inches(top_y),
        Inches(2.3), Inches(top_h),
    )
    ctf = chap_tb.text_frame
    ctf.margin_left = 0
    ctf.margin_right = 0
    ctf.margin_top = 0
    ctf.margin_bottom = 0
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.RIGHT
    cr = cp.add_run()
    cr.text = f"CHAPTER {sec_num_str} / {total_num_str}"
    cr.font.name = FONT_FACE
    cr.font.size = Pt(10)
    cr.font.color.rgb = Color.PRIMARY

    # ─── 巨大な薄色セクション番号 (背景装飾) — スライド左端ぎりぎり ──
    big_num_tb = slide.shapes.add_textbox(
        Inches(0), Inches(2.2),
        Inches(4.5), Inches(3.5),
    )
    btf = big_num_tb.text_frame
    btf.margin_left = 0
    btf.margin_right = 0
    btf.margin_top = 0
    btf.margin_bottom = 0
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.LEFT
    br = bp.add_run()
    br.text = sec_num_str
    br.font.name = FONT_FACE
    br.font.size = Pt(280)
    br.font.bold = True
    br.font.color.rgb = Color.SECTION_BG  # 薄ブルー背景装飾

    # ─── 中央: 小ラベル「基本レポート」+ 大タイトル ──
    label_x = 3.3
    label_y = 3.05
    eyebrow_tb = slide.shapes.add_textbox(
        Inches(label_x), Inches(label_y),
        Inches(SLIDE_W - label_x - 0.7), Inches(0.35),
    )
    etf = eyebrow_tb.text_frame
    etf.margin_left = 0
    etf.margin_right = 0
    ep = etf.paragraphs[0]
    er = ep.add_run()
    er.text = eyebrow
    er.font.name = FONT_FACE
    er.font.size = Pt(13)
    er.font.bold = True
    er.font.color.rgb = Color.PRIMARY

    title_tb = slide.shapes.add_textbox(
        Inches(label_x), Inches(label_y + 0.4),
        Inches(SLIDE_W - label_x - 0.7), Inches(1.3),
    )
    ttf = title_tb.text_frame
    ttf.margin_left = 0
    ttf.margin_right = 0
    tp = ttf.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = FONT_FACE
    tr.font.size = Pt(54)
    tr.font.bold = True
    tr.font.color.rgb = Color.DARK

    # ─── 横線 ──
    line_y = label_y + 1.85
    div_line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(label_x), Inches(line_y),
        Inches(right_x), Inches(line_y),
    )
    div_line.line.color.rgb = Color.BORDER
    div_line.line.width = Pt(0.5)
    _disable_shape_shadow(div_line)

    # ─── 説明文 ──
    desc = SECTION_DESCRIPTIONS.get(title, "")
    if desc:
        desc_tb = slide.shapes.add_textbox(
            Inches(label_x), Inches(line_y + 0.18),
            Inches(SLIDE_W - label_x - 0.7), Inches(1.5),
        )
        dtf = desc_tb.text_frame
        dtf.margin_left = 0
        dtf.margin_right = 0
        dtf.word_wrap = True
        for idx, dline in enumerate(desc.split("\n")):
            dp = dtf.paragraphs[0] if idx == 0 else dtf.add_paragraph()
            dp.line_spacing = 1.6
            dr = dp.add_run()
            dr.text = dline
            dr.font.name = FONT_FACE
            dr.font.size = Pt(13)
            dr.font.color.rgb = Color.SUB_TEXT

    # 共通フッター (© + P.NN)
    _add_slide_footer(slide, ctx)


# ─── チャートヘルパー ────────────────────────────────────


def _apply_line_chart_style(chart, colors_hex: list[str]) -> None:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = FONT_FACE

    for i, series in enumerate(chart.series):
        c = hex_to_rgb(colors_hex[i % len(colors_hex)])
        series.format.line.color.rgb = c
        series.format.line.width = Pt(2)
        # Marker
        try:
            marker = series.marker
            marker.style = 2  # circle
            marker.size = 5
            marker.format.fill.solid()
            marker.format.fill.fore_color.rgb = c
            marker.format.line.color.rgb = c
        except Exception:
            pass

    _style_chart_axes(chart, cat_font=10, val_font=10)


def _apply_column_chart_style(chart, colors_hex: list[str], cat_font: float = 10, rotation: int = 0) -> None:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = FONT_FACE

    for i, series in enumerate(chart.series):
        c = hex_to_rgb(colors_hex[i % len(colors_hex)])
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = c
        series.format.line.fill.background()

    _style_chart_axes(chart, cat_font=cat_font, val_font=10, cat_rotation=rotation)


def _apply_pie_chart_style(chart, colors_hex: list[str], show_percentage: bool = True) -> None:
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(11)
    chart.legend.font.name = FONT_FACE

    # チャートタイトル（系列名の自動タイトル）が巨大化する問題を是正 → 12pt 固定
    try:
        series_name = chart.series[0].name if len(chart.series) else ""
    except Exception:
        series_name = ""
    chart.has_title = True
    tf = chart.chart_title.text_frame
    tf.text = series_name or ""
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(12)
            r.font.name = FONT_FACE
            r.font.bold = True
            r.font.color.rgb = Color.DARK

    plot = chart.plots[0]
    points = plot.series[0].points
    for i, point in enumerate(points):
        try:
            c = hex_to_rgb(colors_hex[i % len(colors_hex)])
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = c
            point.format.line.fill.background()
        except Exception:
            pass

    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_percentage = show_percentage
    dl.show_value = False
    dl.show_category_name = False
    dl.show_legend_key = False
    dl.show_series_name = False
    # NOTE: dLblPos=outEnd はパイ／ドーナツで PowerPoint 互換性問題の報告あり — 既定値を使用
    dl.font.size = Pt(9)
    dl.font.name = FONT_FACE
    dl.number_format = "0.0%"


def _set_plot_area_layout(chart, x: float = 0.08, y: float = 0.05, w: float = 0.88, h: float = 0.75) -> None:
    """<c:plotArea> に manualLayout を追加してプロット領域の相対位置を制御する。

    x,y,w,h はチャート全体に対する比率 (0-1)。既定では高さ 75% に抑えてチャート下端と凡例の間に余白を確保。
    """
    chart_elem = chart._chartSpace
    plot_area = chart_elem.find(".//" + qn("c:plotArea"))
    if plot_area is None:
        return
    # 既存 layout を削除
    for existing in plot_area.findall(qn("c:layout")):
        plot_area.remove(existing)
    layout = etree.Element(qn("c:layout"))
    manual = etree.SubElement(layout, qn("c:manualLayout"))
    for tag, val in (
        ("c:layoutTarget", "inner"),
        ("c:xMode", "edge"),
        ("c:yMode", "edge"),
    ):
        e = etree.SubElement(manual, qn(tag))
        e.set("val", val)
    for tag, val in (
        ("c:x", f"{x}"),
        ("c:y", f"{y}"),
        ("c:w", f"{w}"),
        ("c:h", f"{h}"),
    ):
        e = etree.SubElement(manual, qn(tag))
        e.set("val", f"{val}")
    # plotArea の最初の子として挿入（schema 上 layout が最初）
    plot_area.insert(0, layout)


def _style_chart_axes(chart, cat_font: float = 10, val_font: float = 10, cat_rotation: int = 0) -> None:
    try:
        ca = chart.category_axis
        ca.tick_labels.font.size = Pt(cat_font)
        ca.tick_labels.font.name = FONT_FACE
        if cat_rotation:
            ca.tick_labels.rotation = -abs(cat_rotation)
    except Exception:
        pass
    try:
        va = chart.value_axis
        va.tick_labels.font.size = Pt(val_font)
        va.tick_labels.font.name = FONT_FACE
    except Exception:
        pass


# ─── 3. 月別 ─────────────────────────────────────────────


def _create_monthly_slides(
    prs: Presentation, ctx: _Ctx,
    monthly_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not monthly_sheet or not monthly_sheet.get("rows"):
        return
    rows_desc = monthly_sheet["rows"]  # adapted: 降順
    rows_asc = list(reversed(rows_desc))
    comp_rows_desc = monthly_sheet.get("compRows") or []
    comp_rows_asc = list(reversed(comp_rows_desc))

    # ─── チャートスライド ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide1, Color.WHITE)
    _add_slide_title(slide1, ctx, "月別（13ヶ月推移）")

    labels = [r.get("label", "") for r in rows_asc]
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(short_label_of("totalUsers"), [fmt_num(r.get("users")) for r in rows_asc])
    chart_data.add_series(short_label_of("sessions"), [fmt_num(r.get("sessions")) for r in rows_asc])
    chart_data.add_series(short_label_of("screenPageViews"), [fmt_num(r.get("pageViews")) for r in rows_asc])
    chart_data.add_series(short_label_of("conversions"), [fmt_num(r.get("conversions")) for r in rows_asc])
    colors = ["3b82f6", "f59e0b", "8b5cf6", "ef4444"]
    if comp_rows_asc:
        chart_data.add_series(comparison_label("sessions", "prev", use_short=True), [fmt_num(r.get("sessions")) for r in comp_rows_asc])
        chart_data.add_series(comparison_label("conversions", "prev", use_short=True), [fmt_num(r.get("conversions")) for r in comp_rows_asc])
        colors += ["93c5fd", "fca5a5"]

    chart_h = FOOTER_Y - CONTENT_Y - 0.6
    chart = slide1.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(MARGIN_X), Inches(CONTENT_Y),
        Inches(CONTENT_W), Inches(chart_h),
        chart_data,
    ).chart
    _apply_line_chart_style(chart, colors)

    _add_slide_footer(slide1, ctx)

    # ─── テーブルスライド ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide2, Color.WHITE)
    _add_slide_title(slide2, ctx, "月別（13ヶ月推移）データ")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    headers = [
        {"label": "月", "align": "center"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("totalUsers"), "align": "right"},
        {"label": short_label_of("newUsers"), "align": "right"},
        {"label": short_label_of("screenPageViews"), "align": "right"},
        {"label": short_label_of("engagementRate"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
        {"label": short_label_of("conversionRate"), "align": "right"},
    ]
    rows_out = [
        [
            r.get("label", ""),
            format_number(r.get("sessions")),
            format_number(r.get("users")),
            format_number(r.get("newUsers")),
            format_number(r.get("pageViews")),
            fmt_pct(r.get("engagementRate")),
            format_number(r.get("conversions")),
            fmt_pct(r.get("conversionRate")),
        ]
        for r in rows_desc
    ]
    # 月別は 8 列で横幅統一
    n_cols = len(headers)
    equal_w = [CONTENT_W / n_cols] * n_cols
    _build_table(
        slide2, headers, rows_out, equal_w,
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
        font_size=get_table_font_size(len(rows_out)),
    )
    _add_slide_footer(slide2, ctx)
    _create_ai_slide(prs, ctx, "月別（13ヶ月推移）データ", ai_data, memos)


# ─── 4. 日別 ─────────────────────────────────────────────


def _create_daily_slides(
    prs: Presentation, ctx: _Ctx,
    daily_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not daily_sheet or not daily_sheet.get("rows"):
        return
    rows_desc = daily_sheet["rows"]
    rows_asc = list(reversed(rows_desc))
    comp_rows_desc = daily_sheet.get("compRows") or []
    comp_rows_asc = list(reversed(comp_rows_desc))

    # ─── チャートスライド ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide1, Color.WHITE)
    _add_slide_title(slide1, ctx, "日別推移")

    labels = [_short_date(r.get("date", "")) for r in rows_asc]
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(short_label_of("sessions"), [fmt_num(r.get("sessions")) for r in rows_asc])
    chart_data.add_series(short_label_of("conversions"), [fmt_num(r.get("conversions")) for r in rows_asc])
    colors = ["3b82f6", "ef4444"]
    if comp_rows_asc:
        chart_data.add_series(comparison_label("sessions", "prev", use_short=True), [fmt_num(r.get("sessions")) for r in comp_rows_asc])
        chart_data.add_series(comparison_label("conversions", "prev", use_short=True), [fmt_num(r.get("conversions")) for r in comp_rows_asc])
        colors += ["93c5fd", "fca5a5"]

    chart_h = FOOTER_Y - CONTENT_Y - 0.6
    chart = slide1.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(MARGIN_X), Inches(CONTENT_Y),
        Inches(CONTENT_W), Inches(chart_h),
        chart_data,
    ).chart
    _apply_line_chart_style(chart, colors)
    _style_chart_axes(
        chart,
        cat_font=7 if len(rows_asc) > 31 else 8,
        val_font=10,
        cat_rotation=45 if len(rows_asc) > 15 else 0,
    )
    # 日別推移は項目数が多いため凡例を 6pt に縮小
    chart.legend.font.size = Pt(6)
    _add_slide_footer(slide1, ctx)

    # ─── テーブルスライド（左右 2 列） ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide2, Color.WHITE)
    _add_slide_title(slide2, ctx, "日別推移（データ）")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    headers = [
        {"label": "日付", "align": "center"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
    ]
    rows_out = [
        [r.get("date", ""), format_number(r.get("sessions")), format_number(r.get("conversions"))]
        for r in rows_desc
    ]
    mid = (len(rows_out) + 1) // 2
    left = rows_out[:mid]
    right = rows_out[mid:]
    gap = 0.3
    half_w = (CONTENT_W - gap) / 2
    col_w = [1.6, 1.8, 1.55]
    # スケール
    scale = half_w / sum(col_w)
    col_w_scaled = [w * scale for w in col_w]
    font_size = get_table_font_size(mid)

    _build_table(
        slide2, headers, left, col_w_scaled,
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
        font_size=font_size,
    )
    if right:
        _build_table(
            slide2, headers, right, col_w_scaled,
            x=MARGIN_X + half_w + gap, y=layout["table_y"], h=layout["table_h"],
            font_size=font_size,
        )

    _add_slide_footer(slide2, ctx)
    _create_ai_slide(prs, ctx, "日別推移（データ）", ai_data, memos)


def _short_date(s: str) -> str:
    # "2026/04/01" → "04/01"
    if not s:
        return ""
    parts = str(s).split("/")
    if len(parts) == 3:
        return f"{parts[1]}/{parts[2]}"
    return s


# ─── 5. 曜日別 ───────────────────────────────────────────


def _create_weekly_slide(
    prs: Presentation, ctx: _Ctx,
    weekly_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not weekly_sheet or not weekly_sheet.get("rows"):
        return
    rows = weekly_sheet["rows"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "曜日別")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(7, has_ai)

    labels = [r.get("dayName", "") for r in rows]
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(short_label_of("sessions"), [fmt_num(r.get("sessions")) for r in rows])
    chart_data.add_series(short_label_of("conversions"), [fmt_num(r.get("conversions")) for r in rows])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(MARGIN_X), Inches(layout["chart_y"]),
        Inches(CONTENT_W), Inches(layout["chart_h"]),
        chart_data,
    ).chart
    _apply_column_chart_style(chart, ["3b82f6", "ef4444"], cat_font=8)

    headers = [
        {"label": "曜日", "align": "center"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
    ]
    rows_out = [
        [r.get("dayName", ""), format_number(r.get("sessions")), format_number(r.get("conversions"))]
        for r in rows
    ]
    _build_table(
        slide, headers, rows_out, [3.0, 3.6, 3.5],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
        font_size=10,
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "曜日別", ai_data, memos)


# ─── 6. 時間帯別 ─────────────────────────────────────────


def _create_hourly_slides(
    prs: Presentation, ctx: _Ctx,
    hourly_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not hourly_sheet or not hourly_sheet.get("rows"):
        return
    rows = hourly_sheet["rows"]

    # ─── チャート ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide1, Color.WHITE)
    _add_slide_title(slide1, ctx, "時間帯別")

    labels = [r.get("hour", "") for r in rows]
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series(short_label_of("sessions"), [fmt_num(r.get("sessions")) for r in rows])
    chart_data.add_series(short_label_of("conversions"), [fmt_num(r.get("conversions")) for r in rows])

    chart_h = FOOTER_Y - CONTENT_Y - 0.6
    chart = slide1.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(MARGIN_X), Inches(CONTENT_Y),
        Inches(CONTENT_W), Inches(chart_h),
        chart_data,
    ).chart
    _apply_column_chart_style(chart, ["3b82f6", "ef4444"], cat_font=8)
    _add_slide_footer(slide1, ctx)

    # ─── テーブル（左右 2 列） ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide2, Color.WHITE)
    _add_slide_title(slide2, ctx, "時間帯別（データ）")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    headers = [
        {"label": "時間", "align": "center"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
    ]
    rows_out = [
        [r.get("hour", ""), format_number(r.get("sessions")), format_number(r.get("conversions"))]
        for r in rows
    ]
    mid = (len(rows_out) + 1) // 2
    left = rows_out[:mid]
    right = rows_out[mid:]
    gap = 0.3
    half_w = (CONTENT_W - gap) / 2
    col_w = [1.6, 1.8, 1.55]
    scale = half_w / sum(col_w)
    col_w_scaled = [w * scale for w in col_w]
    font_size = get_table_font_size(mid)

    _build_table(
        slide2, headers, left, col_w_scaled,
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
        font_size=font_size,
    )
    if right:
        _build_table(
            slide2, headers, right, col_w_scaled,
            x=MARGIN_X + half_w + gap, y=layout["table_y"], h=layout["table_h"],
            font_size=font_size,
        )

    _add_slide_footer(slide2, ctx)
    _create_ai_slide(prs, ctx, "時間帯別（データ）", ai_data, memos)


# ─── 7. ユーザー属性（ドーナツ 4） ────────────────────────


def _create_users_donut_slide(
    prs: Presentation, ctx: _Ctx,
    demographics: dict | None,
) -> None:
    if not demographics or not demographics.get("data"):
        return
    data = demographics["data"]
    configs = [
        ("新規/リピーター", data.get("newReturning")),
        ("性別", data.get("gender")),
        ("年齢層", data.get("age")),
        ("デバイス", data.get("device")),
    ]
    if not any(cfg for _, cfg in configs if cfg):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "ユーザー属性")

    # データがあるチャートだけを抽出して詰めて配置（1〜4 本で可変レイアウト）
    valid_configs = [(t, items) for t, items in configs if items]
    n = len(valid_configs)
    if n == 0:
        _add_slide_footer(slide, ctx)
        return

    # 本数に応じたチャートサイズ（縦幅を縮めてタイトル・凡例をドーナツに近づける）
    if n == 1:
        chart_w, chart_h, gap = 4.5, 4.5, 0.0
    elif n == 2:
        chart_w, chart_h, gap = 4.0, 4.5, 0.3
    elif n == 3:
        chart_w, chart_h, gap = 3.0, 4.2, 0.2
    else:  # 4
        chart_w, chart_h, gap = 2.3, 4.0, 0.15

    total_w = chart_w * n + gap * (n - 1)
    x_start = (SLIDE_W - total_w) / 2
    y_chart = 1.5  # 垂直中央付近から配置

    for i, (title, items) in enumerate(valid_configs):
        x = x_start + i * (chart_w + gap)

        labels = [str(d.get("name") or d.get("label") or "") for d in items]
        values = [fmt_num(d.get("value") or d.get("count") or d.get("users") or 0) for d in items]

        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series(title, values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(x), Inches(y_chart),
            Inches(chart_w), Inches(chart_h),
            cd,
        ).chart
        _apply_pie_chart_style(chart, CHART_PALETTE_HEX)
        # チャート内蔵タイトルを使用
        chart.has_title = True
        title_tf = chart.chart_title.text_frame
        title_tf.text = title
        for p in title_tf.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.name = FONT_FACE
                r.font.color.rgb = Color.DARK
        # 凡例は下に
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.legend.include_in_layout = False
        # タイトル直下にドーナツを配置、凡例との余白も最小化
        _set_plot_area_layout(chart, x=0.05, y=0.18, w=0.90, h=0.68)

    _add_slide_footer(slide, ctx)


# ─── 8. ユーザー属性（地域） ─────────────────────────────


def _create_users_region_slide(
    prs: Presentation, ctx: _Ctx,
    demographics: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not demographics or not demographics.get("data"):
        return
    location = demographics["data"].get("location") or {}
    region_data = location.get("region") or location.get("country") or []
    if not region_data:
        return
    top = region_data[:20]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "地域別ランキング Top 20")

    has_ai = bool(ai_data and ai_data.get("summary")) or bool(memos)
    layout = calc_table_only_layout(has_ai)

    headers = [
        {"label": "#", "align": "center"},
        {"label": "地域", "align": "left"},
        {"label": short_label_of("totalUsers"), "align": "right"},
        {"label": "割合", "align": "right"},
    ]
    rows_out = [
        [
            i + 1,
            d.get("name") or "",
            format_number(d.get("value") or d.get("users") or 0),
            f"{float(d.get('percentage')):.1f}%" if d.get("percentage") is not None else "-",
        ]
        for i, d in enumerate(top)
    ]

    # 左右 2 列に分割（上 10 件／下 10 件）
    mid = (len(rows_out) + 1) // 2
    left = rows_out[:mid]
    right = rows_out[mid:]
    gap = 0.3
    half_w = (CONTENT_W - gap) / 2
    col_w = [0.5, 2.8, 1.3, 1.2]
    scale = half_w / sum(col_w)
    col_w_scaled = [w * scale for w in col_w]

    # テーブル高さは 11 行（ヘッダ含む）で 0.32"/行に
    tbl_h = min(layout["table_h"], 0.36 * (mid + 1))
    font_size = get_table_font_size(mid)

    _build_table(
        slide, headers, left, col_w_scaled,
        x=MARGIN_X, y=layout["table_y"], h=tbl_h,
        font_size=font_size,
    )
    if right:
        _build_table(
            slide, headers, right, col_w_scaled,
            x=MARGIN_X + half_w + gap, y=layout["table_y"], h=tbl_h,
            font_size=font_size,
        )

    # AI 分析はテーブル直下に
    ai_y = layout["table_y"] + tbl_h + 0.3
    ai_h = max(FOOTER_Y - ai_y - 0.1, AI_MIN_H)
    _add_ai_and_memo_footer(slide, ai_data, memos, ai_y, ai_h)
    _add_slide_footer(slide, ctx)


# ─── 9. 集客チャネル ─────────────────────────────────────


def _create_channels_slide(
    prs: Presentation, ctx: _Ctx,
    channels_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not channels_sheet or not channels_sheet.get("rows"):
        return
    rows = channels_sheet["rows"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "集客チャネル")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(min(len(rows), 10), has_ai)

    # Top 7 + その他
    sorted_rows = sorted(rows, key=lambda r: -fmt_num(r.get("sessions")))
    top7 = sorted_rows[:7]
    others = sorted_rows[7:]
    pie_labels = [r.get("channelName", "") for r in top7]
    pie_values = [fmt_num(r.get("sessions")) for r in top7]
    others_total = sum(fmt_num(r.get("sessions")) for r in others)
    if others_total > 0:
        pie_labels.append("その他")
        pie_values.append(others_total)

    cd = CategoryChartData()
    cd.categories = pie_labels
    cd.add_series(short_label_of("sessions"), pie_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2.7), Inches(layout["chart_y"]),
        Inches(5.5), Inches(layout["chart_h"]),
        cd,
    ).chart
    _apply_pie_chart_style(chart, CHART_PALETTE_HEX)

    headers = [
        {"label": "チャネル", "align": "left"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("totalUsers"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
    ]
    rows_out = [
        [
            r.get("channelName", ""),
            format_number(r.get("sessions")),
            format_number(r.get("users")),
            format_number(r.get("conversions")),
        ]
        for r in sorted_rows[:10]
    ]
    _build_table(
        slide, headers, rows_out, [3.5, 2.5, 2.2, 1.9],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "集客チャネル", ai_data, memos)


# ─── 10. 流入キーワード ──────────────────────────────────


_FUNNEL_LABELS_JA = {
    "branded": "指名",
    "pureIntent": "純顕在",
    "intent": "顕在",
    "latent": "潜在",
    "noise": "無関係",
}


def _create_keywords_funnel_slide(
    prs: Presentation, ctx: _Ctx,
    keywords_v2: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    """流入キーワード元（V2 ファネル）スライド: 5 層集計 + クラスタ TOP + 改善候補 + AI 分析"""
    if not keywords_v2 or not keywords_v2.get("funnel"):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "流入キーワード元")

    funnel = keywords_v2.get("funnel") or {}
    clusters = keywords_v2.get("clusters") or []
    keywords_list = keywords_v2.get("keywords") or []

    layout = calc_table_only_layout(has_ai=False)

    # 5 層テーブル
    headers = [
        {"label": "層", "align": "left"},
        {"label": "KW数", "align": "right"},
        {"label": "表示", "align": "right"},
        {"label": "クリック", "align": "right"},
        {"label": "CTR", "align": "right"},
        {"label": "平均順位", "align": "right"},
        {"label": "推定 CV", "align": "right"},
        {"label": "主要 KW", "align": "left"},
    ]
    rows_out = []
    for layer_key in ["branded", "pureIntent", "intent", "latent", "noise"]:
        f = funnel.get(layer_key)
        if not f:
            continue
        rows_out.append([
            _FUNNEL_LABELS_JA.get(layer_key, layer_key),
            format_number(f.get("count") or 0),
            format_number(f.get("impressions") or 0),
            format_number(f.get("clicks") or 0),
            f"{float(f.get('ctr') or 0):.2f}%",
            f"{float(f.get('avgPosition') or 0):.1f} 位",
            format_number(f.get("estimatedCV") or 0),
            " / ".join((f.get("topKeywords") or [])[:3]),
        ])

    if rows_out:
        tbl_h = 0.4 * (len(rows_out) + 1)
        # 列幅合計 = CONTENT_W (10.2) に収まるよう調整（旧: 10.4 でスライド右端からはみ出ていた）
        _build_table(
            slide, headers, rows_out, [0.9, 0.7, 1.0, 1.0, 0.8, 1.0, 0.9, 3.9],
            x=MARGIN_X, y=layout["table_y"], h=tbl_h,
            font_size=get_table_font_size(len(rows_out)),
        )

    # クラスタ TOP（あれば）+ 改善候補（あれば）を 2 カラムで下部に
    cluster_top = sorted(clusters, key=lambda c: -(c.get("clicks") or 0))[:3]
    ctr_loss_top = sorted(
        [k for k in keywords_list if k.get("ctrLossFlag") and (k.get("potentialClicks") or 0) > 0],
        key=lambda k: -(k.get("potentialClicks") or 0),
    )[:3]

    bottom_y = layout["table_y"] + 0.4 * (len(rows_out) + 1) + 0.3
    # 残り高さに収まる範囲でテキストボックスを配置
    available_h = max(0.6, FOOTER_Y - bottom_y - 0.1)
    box_h = min(1.6, available_h)
    if (cluster_top or ctr_loss_top) and bottom_y + 0.6 < FOOTER_Y:
        # 簡易テキストで補足表示
        from pptx.util import Inches as _In
        from pptx.util import Pt as _Pt
        # CONTENT_W に収めてスライド右端のオーバーフローを回避（旧: 11" → 0.4" はみ出し）
        tx = slide.shapes.add_textbox(_In(MARGIN_X), _In(bottom_y), _In(CONTENT_W), _In(box_h))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = "意味的クラスタ TOP 3 / 改善候補 TOP 3"
        tf.paragraphs[0].runs[0].font.size = _Pt(12)
        tf.paragraphs[0].runs[0].font.bold = True
        if cluster_top:
            for c in cluster_top:
                p = tf.add_paragraph()
                p.text = f"・[クラスタ] {c.get('name')}: {c.get('keywordCount') or 0} KW / クリック {format_number(c.get('clicks'))} / 中心 {c.get('centerKeyword') or '-'}"
                p.runs[0].font.size = _Pt(10)
        if ctr_loss_top:
            for k in ctr_loss_top:
                p = tf.add_paragraph()
                p.text = f"・[改善候補] {k.get('query')}: 表示 {format_number(k.get('impressions'))} / 順位 {float(k.get('position') or 0):.1f}位 / 潜在 +{format_number(k.get('potentialClicks'))} クリック"
                p.runs[0].font.size = _Pt(10)

    _add_slide_footer(slide, ctx)

    # AI 分析スライドは _create_ai_slide で統一（user_journey と同じ作法）
    _create_ai_slide(prs, ctx, "流入キーワード元", ai_data, memos)


def _create_keywords_slide(
    prs: Presentation, ctx: _Ctx,
    keywords_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not keywords_sheet or not keywords_sheet.get("rows"):
        return
    rows = sorted(keywords_sheet["rows"], key=lambda r: -fmt_num(r.get("clicks")))[:20]

    # ─── スライド1: テーブル（Top20 は長いので AI 分析は別ページへ分離） ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "流入キーワード Top 20")

    layout = calc_table_only_layout(has_ai=False)  # テーブル領域をフル活用

    headers = [
        {"label": "#", "align": "center"},
        {"label": "キーワード", "align": "left"},
        {"label": short_label_of("clicks"), "align": "right"},
        {"label": short_label_of("impressions"), "align": "right"},
        {"label": short_label_of("ctr"), "align": "right"},
        {"label": short_label_of("position"), "align": "right"},
    ]
    rows_out = [
        [
            i + 1,
            r.get("keyword", ""),
            format_number(r.get("clicks")),
            format_number(r.get("impressions")),
            fmt_pct(r.get("ctr")),
            f"{float(r.get('position')):.1f}" if r.get("position") is not None else "-",
        ]
        for i, r in enumerate(rows)
    ]
    _build_table(
        slide, headers, rows_out, [0.6, 4.0, 1.5, 1.5, 1.2, 1.3],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
        font_size=get_table_font_size(len(rows_out)),
    )
    _add_slide_footer(slide, ctx)

    # ─── スライド2: AI 分析 ───
    if (ai_data and ai_data.get("summary")) or memos:
        slide_ai = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_slide_bg(slide_ai, Color.WHITE)
        _add_slide_title(slide_ai, ctx, "流入キーワード Top 20（AI分析）")
        ai_y = CONTENT_Y
        ai_h = FOOTER_Y - CONTENT_Y - 0.1
        _add_ai_and_memo_footer(slide_ai, ai_data, memos, ai_y, ai_h)
        _add_slide_footer(slide_ai, ctx)


# ─── 11. 被リンク元 ──────────────────────────────────────


def _create_referrals_slide(
    prs: Presentation, ctx: _Ctx,
    referrals_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not referrals_sheet or not referrals_sheet.get("rows"):
        return
    rows = referrals_sheet["rows"]
    sorted_rows = sorted(rows, key=lambda r: -fmt_num(r.get("sessions")))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "被リンク元")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(min(len(sorted_rows), 10), has_ai)

    top7 = sorted_rows[:7]
    others = sorted_rows[7:]
    pie_labels = [r.get("source", "") for r in top7]
    pie_values = [fmt_num(r.get("sessions")) for r in top7]
    others_total = sum(fmt_num(r.get("sessions")) for r in others)
    if others_total > 0:
        pie_labels.append("その他")
        pie_values.append(others_total)

    cd = CategoryChartData()
    cd.categories = pie_labels
    cd.add_series(short_label_of("sessions"), pie_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2.7), Inches(layout["chart_y"]),
        Inches(5.5), Inches(layout["chart_h"]),
        cd,
    ).chart
    _apply_pie_chart_style(chart, CHART_PALETTE_HEX)

    headers = [
        {"label": "参照元", "align": "left"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("totalUsers"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
        {"label": short_label_of("engagementRate"), "align": "right"},
    ]
    rows_out = [
        [
            r.get("source", ""),
            format_number(r.get("sessions")),
            format_number(r.get("users")),
            format_number(r.get("conversions")),
            fmt_pct(r.get("engagementRate")),
        ]
        for r in sorted_rows[:10]
    ]
    _build_table(
        slide, headers, rows_out, [3.0, 2.0, 1.8, 1.6, 1.7],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "被リンク元", ai_data, memos)


# ─── 12. ページ別 Top 10 ─────────────────────────────────


def _create_pages_slide(
    prs: Presentation, ctx: _Ctx,
    pages_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not pages_sheet or not pages_sheet.get("rows"):
        return
    sorted_rows = sorted(pages_sheet["rows"], key=lambda r: -fmt_num(r.get("pageViews")))
    top10 = sorted_rows[:10]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "ページ別 Top 10")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(10, has_ai)

    labels = [(r.get("path") or "")[:30] for r in top10]
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series(short_label_of("screenPageViews"), [fmt_num(r.get("pageViews")) for r in top10])
    cd.add_series(short_label_of("sessions"), [fmt_num(r.get("sessions")) for r in top10])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(MARGIN_X), Inches(layout["chart_y"]),
        Inches(CONTENT_W), Inches(layout["chart_h"]),
        cd,
    ).chart
    _apply_column_chart_style(chart, ["3b82f6", "10b981"], cat_font=8, rotation=45)

    headers = [
        {"label": "#", "align": "center"},
        {"label": "ページパス", "align": "left"},
        {"label": short_label_of("screenPageViews"), "align": "right"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("engagementRate"), "align": "right"},
    ]
    rows_out = [
        [
            i + 1,
            (r.get("path") or "")[:50],
            format_number(r.get("pageViews")),
            format_number(r.get("sessions")),
            fmt_pct(r.get("engagementRate")),
        ]
        for i, r in enumerate(top10)
    ]
    _build_table(
        slide, headers, rows_out, [0.6, 4.8, 1.6, 1.6, 1.5],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "ページ別 Top 10", ai_data, memos)


# ─── 13. ページ分類別 ───────────────────────────────────


def _create_page_categories_slide(
    prs: Presentation, ctx: _Ctx,
    page_categories_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not page_categories_sheet or not page_categories_sheet.get("rows"):
        return
    rows = sorted(page_categories_sheet["rows"], key=lambda r: -fmt_num(r.get("pageViews")))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "ページ分類別")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(min(len(rows), 10), has_ai)

    top8 = rows[:8]
    others_total = sum(fmt_num(c.get("pageViews")) for c in rows[8:])
    pie_labels = [c.get("category", "") for c in top8]
    pie_values = [fmt_num(c.get("pageViews")) for c in top8]
    if others_total > 0:
        pie_labels.append("その他")
        pie_values.append(others_total)

    cd = CategoryChartData()
    cd.categories = pie_labels
    cd.add_series(short_label_of("screenPageViews"), pie_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(2.7), Inches(layout["chart_y"]),
        Inches(5.5), Inches(layout["chart_h"]),
        cd,
    ).chart
    _apply_pie_chart_style(chart, CHART_PALETTE_HEX)

    headers = [
        {"label": "カテゴリ", "align": "left"},
        {"label": short_label_of("screenPageViews"), "align": "right"},
        {"label": short_label_of("totalUsers"), "align": "right"},
        {"label": short_label_of("engagementRate"), "align": "right"},
    ]
    rows_out = [
        [
            r.get("category", ""),
            format_number(r.get("pageViews")),
            format_number(r.get("users")),
            fmt_pct(r.get("engagementRate")),
        ]
        for r in rows[:10]
    ]
    _build_table(
        slide, headers, rows_out, [3.5, 2.5, 2.2, 1.9],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "ページ分類別", ai_data, memos)


# ─── 14. ランディングページ ─────────────────────────────


def _create_landing_pages_slide(
    prs: Presentation, ctx: _Ctx,
    landing_sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not landing_sheet or not landing_sheet.get("rows"):
        return
    sorted_rows = sorted(landing_sheet["rows"], key=lambda r: -fmt_num(r.get("sessions")))
    top10 = sorted_rows[:10]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "ランディングページ Top 10")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(10, has_ai)

    labels = [(r.get("path") or "")[:30] for r in top10]
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series(short_label_of("sessions"), [fmt_num(r.get("sessions")) for r in top10])
    cd.add_series(short_label_of("conversions"), [fmt_num(r.get("conversions")) for r in top10])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(MARGIN_X), Inches(layout["chart_y"]),
        Inches(CONTENT_W), Inches(layout["chart_h"]),
        cd,
    ).chart
    _apply_column_chart_style(chart, ["3b82f6", "ef4444"], cat_font=8, rotation=45)

    headers = [
        {"label": "#", "align": "center"},
        {"label": "ランディングページ", "align": "left"},
        {"label": short_label_of("sessions"), "align": "right"},
        {"label": short_label_of("conversions"), "align": "right"},
        {"label": short_label_of("engagementRate"), "align": "right"},
    ]
    rows_out = [
        [
            i + 1,
            (r.get("path") or "")[:50],
            format_number(r.get("sessions")),
            format_number(r.get("conversions")),
            fmt_pct(r.get("engagementRate")),
        ]
        for i, r in enumerate(top10)
    ]
    _build_table(
        slide, headers, rows_out, [0.6, 4.8, 1.6, 1.6, 1.5],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "ランディングページ Top 10", ai_data, memos)


# ─── 15. ファイルDL ─────────────────────────────────────


def _create_file_downloads_slide(
    prs: Presentation, ctx: _Ctx,
    sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not sheet or not sheet.get("rows"):
        return
    sorted_rows = sorted(sheet["rows"], key=lambda r: -fmt_num(r.get("downloads")))
    top10 = sorted_rows[:10]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "ファイルダウンロード")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_layout(min(len(top10), 10), has_ai)

    labels = [(r.get("fileName") or "")[:25] for r in top10]
    cd = CategoryChartData()
    cd.categories = labels
    cd.add_series("DL数", [fmt_num(r.get("downloads")) for r in top10])
    cd.add_series(short_label_of("totalUsers"), [fmt_num(r.get("users")) for r in top10])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(MARGIN_X), Inches(layout["chart_y"]),
        Inches(CONTENT_W), Inches(layout["chart_h"]),
        cd,
    ).chart
    _apply_column_chart_style(chart, ["3b82f6", "10b981"], cat_font=8, rotation=45)

    headers = [
        {"label": "#", "align": "center"},
        {"label": "ファイル名", "align": "left"},
        {"label": "DL数", "align": "right"},
        {"label": short_label_of("totalUsers"), "align": "right"},
    ]
    rows_out = [
        [
            i + 1,
            (r.get("fileName") or "")[:60],
            format_number(r.get("downloads")),
            format_number(r.get("users")),
        ]
        for i, r in enumerate(top10)
    ]
    _build_table(
        slide, headers, rows_out, [0.6, 5.8, 1.9, 1.8],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "ファイルダウンロード", ai_data, memos)


# ─── 16. 外部リンク ─────────────────────────────────────


def _create_external_links_slide(
    prs: Presentation, ctx: _Ctx,
    sheet: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not sheet or not sheet.get("rows"):
        return
    sorted_rows = sorted(sheet["rows"], key=lambda r: -fmt_num(r.get("clicks")))
    top10 = sorted_rows[:10]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "外部リンククリック")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    headers = [
        {"label": "#", "align": "center"},
        {"label": "リンクURL", "align": "left"},
        {"label": "クリック数", "align": "right"},
        {"label": short_label_of("totalUsers"), "align": "right"},
    ]
    rows_out = [
        [
            i + 1,
            (r.get("linkUrl") or "")[:60],
            format_number(r.get("clicks")),
            format_number(r.get("users")),
        ]
        for i, r in enumerate(top10)
    ]
    _build_table(
        slide, headers, rows_out, [0.6, 5.8, 1.9, 1.8],
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "外部リンククリック", ai_data, memos)


# ─── 17. コンバージョン月次推移 ─────────────────────────


def _create_conversions_slides(
    prs: Presentation, ctx: _Ctx,
    conversions: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not conversions or not conversions.get("data"):
        return
    data_asc = sorted(
        conversions["data"],
        key=lambda d: str(d.get("yearMonth") or d.get("month") or ""),
    )
    data_desc = list(reversed(data_asc))

    # CV イベント名（yearMonth / month / total / label 以外のキー）
    excluded = {"yearMonth", "month", "total", "label"}
    event_names: list[str] = []
    for d in data_asc:
        for k in d.keys():
            if k not in excluded and k not in event_names:
                event_names.append(k)

    # ─── チャート ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide1, Color.WHITE)
    _add_slide_title(slide1, ctx, "コンバージョン月次推移")

    labels = [fmt_year_month(d.get("yearMonth") or d.get("month")) for d in data_asc]
    cd = CategoryChartData()
    cd.categories = labels
    colors = ["3b82f6", "10b981", "f59e0b", "ef4444", "8b5cf6", "ec4899", "06b6d4"]
    for ev in event_names:
        cd.add_series(ev, [fmt_num(d.get(ev)) for d in data_asc])
    if event_names:
        chart_h = FOOTER_Y - CONTENT_Y - 0.6
        chart = slide1.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(MARGIN_X), Inches(CONTENT_Y),
            Inches(CONTENT_W), Inches(chart_h),
            cd,
        ).chart
        _apply_line_chart_style(chart, colors)
        # プロット領域の高さを 72% に抑えて、チャート下端と凡例の間に余白を確保
        _set_plot_area_layout(chart, x=0.08, y=0.05, w=0.88, h=0.72)
    _add_slide_footer(slide1, ctx)

    # ─── テーブル ───
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide2, Color.WHITE)
    _add_slide_title(slide2, ctx, "コンバージョン月次推移 データ")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    headers = [{"label": "月", "align": "center"}]
    headers.extend({"label": ev, "align": "right"} for ev in event_names)
    headers.append({"label": "合計", "align": "right"})

    cols_n = len(headers)
    base_w = max((CONTENT_W - 3.0) / max(len(event_names) + 1, 1), 1.0)
    widths = [1.5] + [base_w] * len(event_names) + [1.5]
    scale = CONTENT_W / sum(widths)
    widths = [w * scale for w in widths]

    rows_out = []
    for d in data_desc:
        row = [fmt_year_month(d.get("yearMonth") or d.get("month"))]
        total = 0.0
        for ev in event_names:
            v = fmt_num(d.get(ev))
            total += v
            row.append(format_number(v))
        row.append(format_number(total))
        rows_out.append(row)

    _build_table(
        slide2, headers, rows_out, widths,
        x=MARGIN_X, y=layout["table_y"], h=layout["table_h"],
        font_size=get_table_font_size(len(rows_out)),
    )
    _add_slide_footer(slide2, ctx)
    _create_ai_slide(prs, ctx, "コンバージョン月次推移 データ", ai_data, memos)


# ─── 18. 逆算フロー ─────────────────────────────────────


def _create_reverse_flow_slide(
    prs: Presentation, ctx: _Ctx,
    reverse_flows: list | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    if not reverse_flows:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "逆算フロー")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    headers = [
        {"label": "フロー名", "align": "left"},
        {"label": "全PV", "align": "right"},
        {"label": "遷移率①", "align": "right"},
        {"label": "フォームPV", "align": "right"},
        {"label": "遷移率②", "align": "right"},
        {"label": "送信完了", "align": "right"},
        {"label": "全体CVR", "align": "right"},
    ]
    rows_out = []
    for flow in reverse_flows:
        s = flow.get("summary") or {}
        total_pv = fmt_num(s.get("totalSiteViews"))
        form_pv = fmt_num(s.get("formPageViews"))
        complete = fmt_num(s.get("submissionComplete"))
        rate1 = f"{(form_pv / total_pv * 100):.2f}%" if total_pv > 0 else "-"
        rate2 = f"{(complete / form_pv * 100):.2f}%" if form_pv > 0 else "-"
        cvr = f"{(complete / total_pv * 100):.2f}%" if total_pv > 0 else "-"
        rows_out.append([
            flow.get("flowName") or "",
            format_number(total_pv),
            rate1,
            format_number(form_pv),
            rate2,
            format_number(complete),
            cvr,
        ])
    # テーブルは行数に応じた高さで、その直下に AI 分析を配置
    tbl_h = 0.4 * (len(rows_out) + 1)
    _build_table(
        slide, headers, rows_out, [2.5, 1.3, 1.2, 1.3, 1.2, 1.3, 1.3],
        x=MARGIN_X, y=layout["table_y"], h=tbl_h,
    )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "逆算フロー", ai_data, memos)


# ─── 17.5 ユーザージャーニー ────────────────────────────


def _create_user_journey_slide(
    prs: Presentation, ctx: _Ctx,
    user_journey: dict | None,
    ai_data: dict | None,
    memos: list | None,
) -> None:
    """ユーザージャーニー（5層フロー）スライド: 主要ジャーニー TOP 3 + 詳細パス TOP 6"""
    if not user_journey or not user_journey.get("nodes"):
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "ユーザージャーニー")

    has_ai = bool(ai_data and ai_data.get("summary"))
    layout = calc_table_only_layout(has_ai)

    detail_paths = user_journey.get("detailPaths") or []
    headers = [
        {"label": "#", "align": "right"},
        {"label": "流入元", "align": "left"},
        {"label": "ランディング", "align": "left"},
        {"label": "中間", "align": "left"},
        {"label": "結果", "align": "left"},
        {"label": "セッション", "align": "right"},
        {"label": "CV 率", "align": "right"},
    ]
    rows_out = []
    for p in detail_paths[:8]:
        if not isinstance(p, dict):
            continue
        rows_out.append([
            str(p.get("rank") or ""),
            p.get("source") or "",
            p.get("lp") or "",
            p.get("middle") or "—",
            p.get("result") or "",
            format_number(p.get("sessions") or 0),
            f"{float(p.get('cvRate') or 0):.1f}%",
        ])

    if rows_out:
        tbl_h = 0.4 * (len(rows_out) + 1)
        # 結果列 (CV イベント名 14+ 文字) を広め、その分 流入元/セッション を圧縮。合計 = CONTENT_W (10.2)
        _build_table(
            slide, headers, rows_out, [0.5, 1.3, 3.0, 1.8, 1.5, 1.1, 1.0],
            x=MARGIN_X, y=layout["table_y"], h=tbl_h,
            font_size=get_table_font_size(len(rows_out)),
        )
    _add_slide_footer(slide, ctx)
    _create_ai_slide(prs, ctx, "ユーザージャーニー", ai_data, memos)


# ─── 19. Appendix 用語集 ────────────────────────────────


# 用語集に載せる指標 (辞書の canonical key)。ランディングページは辞書外なので
# 固定テキストのタプルとして末尾に追加する。
_APPENDIX_METRIC_KEYS = [
    "sessions",
    "totalUsers",
    "newUsers",
    "screenPageViews",
    "engagementRate",
    "averageSessionDuration",
    "conversions",
    "conversionRate",
]

_APPENDIX_EXTRA_TERMS = [
    ("ランディングページ", "ユーザーが最初に閲覧したページ"),
]

_APPENDIX_GSC_KEYS = ["clicks", "impressions", "ctr", "position"]


def _build_appendix_terms() -> list[tuple[str, str]]:
    terms: list[tuple[str, str]] = []
    for k in _APPENDIX_METRIC_KEYS:
        terms.append((label_of(k), tooltip_of(k)))
    terms.extend(_APPENDIX_EXTRA_TERMS)
    for k in _APPENDIX_GSC_KEYS:
        terms.append((label_of(k), tooltip_of(k)))
    return terms


APPENDIX_TERMS = _build_appendix_terms()


# ─── 改善提案アクションプラン (1 スライド一覧) ────────────────

import re as _re

_IMPR_MOCKUP_SHARE_BASE_URL = "https://grow-reporter.com"
_IMPR_CATEGORY_LABELS = {
    "acquisition": "集客",
    "content": "コンテンツ",
    "design": "デザイン",
    "feature": "機能",
    "other": "その他",
}
_IMPR_PRIORITY_LABELS = {"high": "高", "medium": "中", "low": "低"}


def _impr_build_mockup_share_url(mockup_storage_url):
    if not mockup_storage_url:
        return ""
    m = _re.search(r"/page-mockups/([^/]+)/([^/]+)\.html", str(mockup_storage_url))
    if not m:
        return ""
    return f"{_IMPR_MOCKUP_SHARE_BASE_URL}/page-mockups/{m.group(1)}/{m.group(2)}.html"


def _impr_get_timestamp(ts) -> float:
    if ts is None:
        return 0
    if isinstance(ts, (int, float)):
        return float(ts)
    if isinstance(ts, dict):
        return float(ts.get("seconds") or ts.get("_seconds") or 0)
    return 0


def _create_improvements_slide(prs: Presentation, ctx: _Ctx, improvements: list | None) -> None:
    """改善提案アクションプラン (複数スライド分割、1 スライド最大 10 件)。
    - URL あり: 「モックアップを開く」(ハイパーリンク)
    - mockupSkipped: 「対応不要」
    - 通常: 「未生成」
    """
    if not improvements:
        return

    sorted_items = sorted(
        improvements,
        key=lambda x: (
            x.get("order") if x.get("order") is not None else 999999,
            _impr_get_timestamp(x.get("createdAt")),
        ),
    )

    PER_SLIDE = 10
    total = len(sorted_items)
    total_pages = (total + PER_SLIDE - 1) // PER_SLIDE

    # タイトル列に期待効果を 2 段目に積むレイアウト (期待効果列は削除して横幅圧縮)
    # CONTENT_W = SLIDE_W (11.0) - 2 * MARGIN_X (0.4) = 10.2
    headers = [
        {"label": "No.", "align": "center"},
        {"label": "カテゴリ", "align": "center"},
        {"label": "優先度", "align": "center"},
        {"label": "タイトル / 期待効果", "align": "left"},
        {"label": "モックアップ", "align": "center"},
    ]
    col_widths = [0.35, 0.95, 0.55, 5.6, 2.75]  # 合計 = 10.2
    TITLE_COL = 3
    MOCKUP_COL = 4

    # 全ページで行高さを統一するため、PER_SLIDE 件想定の row 高さを基準にする
    # _build_table は h / rows_n で行高さを決めるため、少ない行のページでは
    # 個別に小さい h を渡すことで巨大化を防ぐ
    table_h_full = FOOTER_Y - CONTENT_Y - 0.15
    uniform_row_h = table_h_full / (PER_SLIDE + 1)  # +1 はヘッダー行

    for page_idx in range(total_pages):
        start = page_idx * PER_SLIDE
        end = min(start + PER_SLIDE, total)
        chunk = sorted_items[start:end]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_slide_bg(slide, Color.WHITE)
        # ページ番号付きタイトル (複数スライドの場合のみ)
        title_label = "改善提案アクションプラン"
        if total_pages > 1:
            title_label = f"改善提案アクションプラン（{page_idx + 1}/{total_pages}）"
        _add_slide_title(slide, ctx, title_label)

        table_y = CONTENT_Y
        # 行高さを統一: このスライドの行数 + ヘッダー = (n + 1) * uniform_row_h
        table_h = uniform_row_h * (len(chunk) + 1)

        # 表示用テキストと URL を事前計算
        mockup_cells: list[tuple[str, str | None]] = []
        for item in chunk:
            url = _impr_build_mockup_share_url(item.get("mockupStorageUrl"))
            if url:
                mockup_cells.append(("モックアップを開く", url))
            elif item.get("mockupSkipped"):
                mockup_cells.append(("対応不要", None))
            else:
                mockup_cells.append(("未生成", None))

        rows_data = []
        for i, item in enumerate(chunk):
            cat = item.get("category") or ""
            pri = item.get("priority") or ""
            # タイトル列はあとで 2 段組に上書きするため、いったんタイトルだけ入れる
            rows_data.append([
                str(start + i + 1),
                _IMPR_CATEGORY_LABELS.get(cat, cat),
                _IMPR_PRIORITY_LABELS.get(pri, pri),
                item.get("title") or "",
                mockup_cells[i][0],
            ])

        _build_table(
            slide, headers, rows_data, col_widths,
            x=MARGIN_X, y=table_y, h=table_h,
            font_size=get_table_font_size(len(rows_data)),
        )

        last_shape = slide.shapes[-1]
        if last_shape.has_table:
            table = last_shape.table
            font_size_data = get_table_font_size(len(rows_data))
            # 1) タイトル列を「タイトル + 期待効果」の 2 段組に上書き
            for i, item in enumerate(chunk):
                cell = table.cell(i + 1, TITLE_COL)
                title_text = item.get("title") or ""
                impact_text = item.get("expectedImpact") or ""
                # クリアして 2 パラグラフ追加
                cell.text = ""
                tf = cell.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                # 1 段目: タイトル
                p1 = tf.paragraphs[0]
                p1.alignment = PP_ALIGN.LEFT
                r1 = p1.add_run()
                r1.text = title_text
                r1.font.name = FONT_FACE
                r1.font.size = Pt(font_size_data)
                r1.font.color.rgb = Color.DARK
                # 2 段目: 期待効果 (グレー、小さめ)
                if impact_text:
                    p2 = tf.add_paragraph()
                    p2.alignment = PP_ALIGN.LEFT
                    p2.line_spacing = 1.3
                    r2 = p2.add_run()
                    r2.text = impact_text
                    r2.font.name = FONT_FACE
                    r2.font.size = Pt(max(font_size_data - 1.5, 7.5))
                    r2.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)  # gray-500

            # 2) モックアップ列の URL 行をハイパーリンク化
            for r, (_text, url) in enumerate(mockup_cells, start=1):
                if not url:
                    continue
                cell = table.cell(r, MOCKUP_COL)
                p = cell.text_frame.paragraphs[0]
                if p.runs:
                    run = p.runs[0]
                    run.hyperlink.address = url
                    run.font.color.rgb = RGBColor(0x05, 0x63, 0xC1)
                    run.font.underline = True

        _add_slide_footer(slide, ctx)


def _create_appendix_slide(prs: Presentation, ctx: _Ctx) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_slide_bg(slide, Color.WHITE)
    _add_slide_title(slide, ctx, "指標・用語の説明")

    table_y = CONTENT_Y
    table_h = FOOTER_Y - table_y - 0.15

    headers = [
        {"label": "指標名", "align": "left"},
        {"label": "説明", "align": "left"},
    ]
    rows_out = [[term, desc] for term, desc in APPENDIX_TERMS]
    _build_table(
        slide, headers, rows_out, [3.0, 7.1],
        x=MARGIN_X, y=table_y, h=table_h,
        font_size=9,
    )
    _add_slide_footer(slide, ctx)
